import base64
import hashlib
import json
import mimetypes
import re
import uuid
import shutil
import subprocess
import sys
import tempfile
from datetime import datetime
from pathlib import Path
from urllib.parse import unquote

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from PIL import Image
from openai import OpenAI
try:
    from utils.token_counter import get_token_count
except ModuleNotFoundError:
    from token_counter import get_token_count

from models import Document
from utils.ai_endpoint import get_ai_base_url_alt
from utils.error_codes import BizCode
from utils.logo_only_filter import LogoOnlyFilter
from utils.ppt_template_cleaner import clean_pptx_template
from utils.ppt_noise_filter import PPTNoiseFilter
from utils.title_utils import normalize_document_title
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

"""
解析ppt，使用python-pptx提取ppt中的图像和文本
"""

def _env_bool(name: str, default: bool = False) -> bool:
    value = os.getenv(name)
    if value is None:
        return default
    return value.strip().lower() in {"1", "true", "yes", "on", "y"}


def _env_float(name: str, default: float) -> float:
    try:
        return float(os.getenv(name, str(default)))
    except (TypeError, ValueError):
        return default


def _env_int(name: str, default: int) -> int:
    try:
        return int(os.getenv(name, str(default)))
    except (TypeError, ValueError):
        return default


class PPTParser:
    NATIVE_TEXT_MIN_CHARS = 80
    IMAGE_AREA_RATIO = 0.6
    DOCUMENT_ROUTE_RATIO = 0.7

    def __init__(self):
        self.document_base_dir = os.getenv("DOCUMENT_BASE_DIR", "D:/Pycharm/code/Maintenance_Assistance_System")
        self.document_dir = os.getenv("DOCUMENT_DIR", "upload/documents")
        self.image_dir = os.getenv("IMAGE_DIR", "upload/images")
        self.api_key = os.getenv("API_KEY", "EMPTY")
        self.model = os.getenv("MODEL_AI", "/models/Qwen3-VL-8B-Instruct")
        self.max_token = int(os.getenv("MAX_TOKEN", 2000))
        self.input_token = int(os.getenv("INPUT_TOKEN", 8000))
        self.model_image_max_size = int(os.getenv("MODEL_IMAGE_MAX_SIZE", 1024))
        self.ppt_native_text_min_chars = self.NATIVE_TEXT_MIN_CHARS
        self.ppt_image_area_ratio = self.IMAGE_AREA_RATIO
        self.ppt_document_route_ratio = self.DOCUMENT_ROUTE_RATIO
        self.ppt_to_pdf_enabled = _env_bool("PPT_TO_PDF_ENABLED", True)
        self.ppt_include_page_images = _env_bool("PPT_INCLUDE_PAGE_IMAGES", False)
        self.ppt_skip_decorative_images = _env_bool("PPT_SKIP_DECORATIVE_IMAGES", True)
        self.ppt_min_image_area_ratio = _env_float("PPT_MIN_IMAGE_AREA_RATIO", 0.015)
        self.ppt_background_image_area_ratio = _env_float("PPT_BACKGROUND_IMAGE_AREA_RATIO", 0.65)
        self.ppt_repeated_image_slide_threshold = _env_int("PPT_REPEATED_IMAGE_SLIDE_THRESHOLD", 2)
        self.ppt_decorative_text_min_chars = _env_int("PPT_DECORATIVE_TEXT_MIN_CHARS", 30)
        self.ppt_soffice_exe = os.getenv("PPT_SOFFICE_EXE", "").strip()
        self.ppt_keep_converted_pdf = _env_bool("PPT_KEEP_CONVERTED_PDF", False)
        self.ppt_clean_template_enabled = _env_bool("PPT_CLEAN_TEMPLATE_ENABLED", True)
        self.ppt_keep_cleaned_ppt = _env_bool("PPT_KEEP_CLEANED_PPT", False)
        self.ppt_clean_slide_placeholder_noise = _env_bool("PPT_CLEAN_SLIDE_PLACEHOLDER_NOISE", True)
        self.ppt_parse_mode = os.getenv("PPT_PARSE_MODE", "strategy").strip().lower() or "strategy"
        if self.ppt_parse_mode not in {"mineru", "mineru_first", "pdf_mineru", "always", "auto", "strategy", "native"}:
            self.ppt_parse_mode = "strategy"
        self.mineru_exe = os.getenv("MINERU_EXE", "").strip()
        self.mineru_timeout = _env_int("MINERU_TIMEOUT", 1800)
        self.ppt_repeated_image_max_area_ratio = _env_float("PPT_REPEATED_IMAGE_MAX_AREA_RATIO", 0.08)
        self.logo_only_filter = LogoOnlyFilter()
        self.ppt_noise_filter = PPTNoiseFilter()
        self.last_parse_strategy = None
        self.last_template_cleaning = None
        self.last_error_code = None
        self.last_error_detail = None
        base_url = os.path.join(self.document_base_dir, self.document_dir)
        if not os.path.exists(base_url):
            os.makedirs(base_url)
            print(f"创建目录: {base_url}")
        base_url = os.path.join(self.document_base_dir, self.image_dir)
        if not os.path.exists(base_url):
            os.makedirs(base_url)
            print(f"创建目录: {base_url}")

    def parse(self, file_path: str):
        self.last_error_code = None
        self.last_error_detail = None
        self.last_template_cleaning = None
        cleanup_context = None
        parse_file_path = file_path
        try:
            parse_file_path, cleanup_context, clean_stats = self._prepare_ppt_source(file_path)
            self.last_template_cleaning = clean_stats
            self.last_parse_strategy = {
                "strategy": "native" if self.ppt_parse_mode == "native" else "pdf_mineru_first",
                "reason": "PPT 默认先转 PDF 后使用 MinerU 解析 Markdown。",
            }
            if clean_stats:
                self.last_parse_strategy["template_cleaning"] = clean_stats

            if self.ppt_parse_mode in {"mineru", "mineru_first", "pdf_mineru", "always"}:
                document = self._parse_with_mineru_if_available(parse_file_path)
                if document is not None:
                    return document

            try:
                strategy = self.get_parse_strategy(parse_file_path)
                self.last_parse_strategy = strategy
                if clean_stats:
                    self.last_parse_strategy["template_cleaning"] = clean_stats
                print(
                    "[PPTParser] fallback strategy={strategy} image_page_ratio={image_ratio:.2f} "
                    "native_page_ratio={native_ratio:.2f} slides={slides}".format(
                        strategy=strategy["strategy"],
                        image_ratio=strategy["image_page_ratio"],
                        native_ratio=strategy["native_page_ratio"],
                        slides=strategy["slide_count"],
                    )
                )
            except Exception as error:
                print(f"[PPTParser] 原生 PPT 策略分析失败，继续尝试原生解析：{error}")

            if self._should_try_mineru_route(self.last_parse_strategy):
                if self.last_parse_strategy is not None:
                    self.last_parse_strategy["selected_route"] = "mineru"
                document = self._parse_with_mineru_if_available(parse_file_path)
                if document is not None:
                    return document

            if self.last_parse_strategy is not None:
                self.last_parse_strategy["selected_route"] = "native"

            text, image_urls, image_names, section_image_indexes = self.get_content(parse_file_path)
            document = self.file2document(text, image_urls, image_names, section_image_indexes)
            return document
        finally:
            if cleanup_context is not None:
                cleanup_context.cleanup()

    def _prepare_ppt_source(self, file_path: str):
        if not self.ppt_clean_template_enabled:
            return file_path, None, None

        source_path = Path(file_path)
        if source_path.suffix.lower() != ".pptx":
            return file_path, None, None

        cleanup_context = None
        try:
            runtime_root = Path(self.document_base_dir) / "runtime" / "ppt_cleaned"
            runtime_root.mkdir(parents=True, exist_ok=True)

            if self.ppt_keep_cleaned_ppt:
                output_dir = runtime_root / f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:8]}"
                output_dir.mkdir(parents=True, exist_ok=True)
            else:
                cleanup_context = tempfile.TemporaryDirectory(prefix="ppt_clean_", dir=str(runtime_root))
                output_dir = Path(cleanup_context.name)

            cleaned_path = output_dir / source_path.name
            clean_stats = clean_pptx_template(
                str(source_path),
                str(cleaned_path),
                remove_slide_placeholder_noise=self.ppt_clean_slide_placeholder_noise,
                remove_slide_noise=True,
                noise_filter=self.ppt_noise_filter,
            )
            print(
                "[PPTParser] 已先清模板：templates={templates} slide_placeholders={placeholders} backgrounds={backgrounds}".format(
                    templates=clean_stats.get("removed_template_shapes", 0),
                    placeholders=clean_stats.get("removed_slide_placeholders", 0),
                    backgrounds=clean_stats.get("removed_backgrounds", 0),
                )
            )
            print(
                "[PPTParser] slide noise filtered: shapes={shapes} text={text} pictures={pictures}".format(
                    shapes=clean_stats.get("removed_slide_noise_shapes", 0),
                    text=clean_stats.get("removed_slide_noise_text_shapes", 0),
                    pictures=clean_stats.get("removed_slide_noise_picture_shapes", 0),
                )
            )
            return str(cleaned_path), cleanup_context, clean_stats
        except Exception as error:
            if cleanup_context is not None:
                cleanup_context.cleanup()
            print(f"[PPTParser] 清模板失败，继续使用原始 PPT：{error}")
            return file_path, None, None

    def _set_last_error(self, code: int, message: str):
        self.last_error_code = int(code)
        self.last_error_detail = message

    def _is_ai_service_unavailable_error(self, error: Exception) -> bool:
        name = type(error).__name__
        if name in {"APIConnectionError", "APITimeoutError"}:
            return True
        msg = str(error).lower()
        keywords = [
            "connection",
            "timed out",
            "timeout",
            "refused",
            "temporarily unavailable",
            "service unavailable",
            "name resolution",
            "max retries exceeded",
            "502",
            "503",
            "504",
        ]
        return any(k in msg for k in keywords)

    def _detect_section(self, text: str):
        normalized = "".join((text or "").split())
        section_keywords = [
            ("problem_intro", ["问题描述", "问题简介", "现场问题描述", "客户反馈", "问题现象", "调查报告"]),
            ("causes", ["原因分析", "推测原因", "根因", "原因"]),
            ("evaluation", ["故障评估", "评估", "影响范围", "风险评估"]),
            ("inspection", ["检查步骤", "检查", "排查记录", "排查", "测试记录", "验证记录"]),
            ("solutions", ["解决方案", "问题解决", "解决", "改善措施", "纠正措施", "建议"]),
            ("key_points", ["关键要点", "总结", "结论"]),
        ]
        for section, keywords in section_keywords:
            if any(keyword in normalized for keyword in keywords):
                return section
        return None

    def _empty_section_image_indexes(self):
        return {
            "problem_intro": [],
            "causes": [],
            "evaluation": [],
            "inspection": [],
            "solutions": [],
            "key_points": [],
        }

    def _iter_shapes(self, shapes):
        for shape in shapes:
            yield shape
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                yield from self._iter_shapes(shape.shapes)

    def _shape_text(self, shape) -> str:
        parts = []
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    parts.append(text)
        if getattr(shape, "has_table", False):
            table_text = self._table_to_markdown(shape.table)
            if table_text:
                parts.append(table_text)
        return self.ppt_noise_filter.sanitize_text("\n".join(parts))

    def _shape_area_ratio(self, shape, slide_width: int, slide_height: int) -> float:
        slide_area = max(int(slide_width) * int(slide_height), 1)
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        return min(max(width * height, 0) / slide_area, 1.0)

    def _is_picture_shape(self, shape) -> bool:
        return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE

    def _picture_hash(self, shape) -> str:
        try:
            return hashlib.sha1(shape.image.blob).hexdigest()
        except Exception:
            return ""

    def _collect_picture_hash_counts(self, prs) -> dict:
        counts = {}
        if not self.ppt_skip_decorative_images:
            return counts
        for slide in prs.slides:
            slide_hashes = set()
            for shape in self._iter_shapes(slide.shapes):
                if self._is_picture_shape(shape):
                    picture_hash = self._picture_hash(shape)
                    if picture_hash:
                        slide_hashes.add(picture_hash)
            for picture_hash in slide_hashes:
                counts[picture_hash] = counts.get(picture_hash, 0) + 1
        return counts

    def _should_keep_picture_shape(
        self,
        shape,
        slide_analysis: dict,
        slide_width: int,
        slide_height: int,
        picture_hash_counts: dict,
    ) -> bool:
        try:
            if self.logo_only_filter.should_filter_bytes(
                shape.image.blob,
                source=f"ppt-shape:{getattr(shape, 'name', 'picture')}",
            ):
                return False
        except Exception:
            pass

        area_ratio = self._shape_area_ratio(shape, slide_width, slide_height)
        try:
            should_filter, reason = self.ppt_noise_filter.should_filter_image_bytes(
                shape.image.blob,
                area_ratio=area_ratio,
                source=f"ppt-shape:{getattr(shape, 'name', 'picture')}",
            )
            if should_filter:
                print(
                    f"[PPTNoiseFilter] filtered PPT image reason={reason} "
                    f"shape={getattr(shape, 'name', 'picture')}",
                    flush=True,
                )
                return False
        except Exception:
            pass

        if not self.ppt_skip_decorative_images:
            return True

        if area_ratio < self.ppt_min_image_area_ratio:
            return False

        picture_hash = self._picture_hash(shape)
        if (
            picture_hash
            and picture_hash_counts.get(picture_hash, 0) >= self.ppt_repeated_image_slide_threshold
            and area_ratio <= self.ppt_repeated_image_max_area_ratio
        ):
            return False

        has_native_content = (
            int(slide_analysis.get("native_chars", 0) or 0) >= self.ppt_decorative_text_min_chars
            or int(slide_analysis.get("table_count", 0) or 0) > 0
        )
        if has_native_content and area_ratio >= self.ppt_background_image_area_ratio:
            return False

        return True

    def _table_to_markdown(self, table) -> str:
        rows = []
        for row in table.rows:
            values = []
            for cell in row.cells:
                value = " ".join((cell.text or "").split())
                values.append(value.replace("|", "\\|"))
            if any(values):
                rows.append(values)
        if not rows:
            return ""
        max_cols = max(len(row) for row in rows)
        normalized_rows = [row + [""] * (max_cols - len(row)) for row in rows]
        header = normalized_rows[0]
        separator = ["---"] * max_cols
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(separator) + " |",
        ]
        for row in normalized_rows[1:]:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)

    def _analyze_slide(self, slide, slide_index: int, slide_width: int, slide_height: int):
        native_chars = 0
        text_shape_count = 0
        table_count = 0
        picture_count = 0
        max_picture_area_ratio = 0.0

        for shape in self._iter_shapes(slide.shapes):
            text = self._shape_text(shape)
            if text:
                text_shape_count += 1
                native_chars += len(text)
            if getattr(shape, "has_table", False):
                table_count += 1
            if self._is_picture_shape(shape):
                picture_count += 1
                max_picture_area_ratio = max(
                    max_picture_area_ratio,
                    self._shape_area_ratio(shape, slide_width, slide_height),
                )

        is_image_page = (
            native_chars < self.ppt_native_text_min_chars
            and table_count == 0
            and picture_count > 0
            and max_picture_area_ratio >= self.ppt_image_area_ratio
        )
        is_native_page = native_chars >= self.ppt_native_text_min_chars or table_count > 0
        return {
            "slide_index": slide_index,
            "native_chars": native_chars,
            "text_shape_count": text_shape_count,
            "table_count": table_count,
            "picture_count": picture_count,
            "max_picture_area_ratio": round(max_picture_area_ratio, 4),
            "is_image_page": is_image_page,
            "is_native_page": is_native_page,
        }

    def get_parse_strategy(self, file_path: str):
        if not os.path.exists(file_path):
            raise FileNotFoundError(file_path)
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        slide_stats = [
            self._analyze_slide(
                slide,
                slide_index,
                int(prs.slide_width),
                int(prs.slide_height),
            )
            for slide_index, slide in enumerate(prs.slides, start=1)
        ]
        image_page_count = sum(1 for item in slide_stats if item["is_image_page"])
        native_page_count = sum(1 for item in slide_stats if item["is_native_page"])
        denominator = max(slide_count, 1)
        image_page_ratio = image_page_count / denominator
        native_page_ratio = native_page_count / denominator

        if image_page_ratio >= self.ppt_document_route_ratio:
            strategy = "mineru"
            reason = "图片型页面占比达到阈值，适合转 PDF 后用 MinerU OCR"
        elif native_page_ratio >= self.ppt_document_route_ratio:
            strategy = "native"
            reason = "原生文本/表格页面占比达到阈值，适合直接读取 PPT 结构"
        else:
            strategy = "mixed"
            reason = "原生页面和图片型页面混合，优先尝试 MinerU，失败后用原生结构兜底"

        return {
            "strategy": strategy,
            "reason": reason,
            "slide_count": slide_count,
            "image_page_count": image_page_count,
            "native_page_count": native_page_count,
            "image_page_ratio": image_page_ratio,
            "native_page_ratio": native_page_ratio,
            "thresholds": {
                "native_text_min_chars": self.ppt_native_text_min_chars,
                "image_area_ratio": self.ppt_image_area_ratio,
                "document_route_ratio": self.ppt_document_route_ratio,
            },
            "slides": slide_stats,
        }

    def _should_try_mineru_route(self, strategy: dict | None) -> bool:
        if self.ppt_parse_mode in {"native", "mineru", "mineru_first", "pdf_mineru", "always"}:
            return False
        if not strategy:
            return False
        if strategy.get("strategy") == "mineru":
            return True
        if strategy.get("strategy") == "mixed":
            return int(strategy.get("image_page_count") or 0) > 0 and int(strategy.get("native_page_count") or 0) == 0
        return False

    def _resolve_mineru_executable(self) -> str:
        if self.mineru_exe:
            mineru_path = os.path.abspath(os.path.expanduser(self.mineru_exe))
            if os.path.exists(mineru_path):
                return mineru_path
            raise FileNotFoundError(f"MINERU_EXE 指定的文件不存在：{mineru_path}")

        python_exe = Path(sys.executable)
        candidates = []
        if os.name == "nt":
            candidates.extend([
                python_exe.parent / "mineru.exe",
                python_exe.parent / "mineru.cmd",
                python_exe.parent / "mineru.bat",
            ])
        else:
            candidates.append(python_exe.parent / "mineru")

        for candidate in candidates:
            if candidate.exists():
                return str(candidate)

        path_command = shutil.which("mineru")
        if path_command:
            return path_command

        raise FileNotFoundError(
            "没有找到 MinerU 可执行文件。请确认已在当前虚拟环境安装 MinerU，"
            "或者在 .env 中设置 MINERU_EXE。"
        )

    def _find_mineru_markdown(self, output_dir: str) -> Path:
        markdown_files = list(Path(output_dir).rglob("*.md"))
        if not markdown_files:
            raise FileNotFoundError(
                f"MinerU 已执行，但在输出目录中没有找到 .md 文件：{output_dir}"
            )
        return max(
            markdown_files,
            key=lambda path: (path.stat().st_size, path.stat().st_mtime),
        )

    def _copy_mineru_images(self, markdown_text: str, markdown_dir: Path):
        image_urls = []
        image_names = []
        image_indexes = {}
        target_dir = Path(self.document_base_dir) / self.image_dir
        target_dir.mkdir(parents=True, exist_ok=True)
        markdown_root = markdown_dir.resolve()

        def copy_image(raw_path: str) -> int | None:
            relative_path = unquote(raw_path.strip().strip('"\''))
            relative_path = relative_path.split("#", 1)[0].split("?", 1)[0]
            source_path = (markdown_dir / relative_path).resolve()

            try:
                source_path.relative_to(markdown_root)
            except ValueError:
                return None
            if not source_path.is_file():
                return None

            source_key = str(source_path).lower()
            if source_key in image_indexes:
                return image_indexes[source_key]
            if self.logo_only_filter.should_filter_path(source_path):
                return None
            should_filter, reason = self.ppt_noise_filter.should_filter_image_path(
                source_path,
                source="mineru-ppt-image",
            )
            if should_filter:
                print(
                    f"[PPTNoiseFilter] filtered MinerU image reason={reason} "
                    f"path={source_path.name}",
                    flush=True,
                )
                image_indexes[source_key] = None
                return None

            extension = source_path.suffix.lower() or ".png"
            unique_filename = f"{uuid.uuid4().hex}{extension}"
            target_path = target_dir / unique_filename
            shutil.copy2(source_path, target_path)

            image_urls.append(str(target_path))
            image_names.append(unique_filename)
            image_index = len(image_urls)
            image_indexes[source_key] = image_index
            print(f"[MinerU] 已保存 PPT 图片：{unique_filename}")
            return image_index

        markdown_image_pattern = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")

        def replace_markdown_image(match: re.Match) -> str:
            image_index = copy_image(match.group(1))
            return f"\n【图片{image_index}】\n" if image_index else "\n"

        markdown_text = markdown_image_pattern.sub(replace_markdown_image, markdown_text)

        html_image_pattern = re.compile(
            r"<img\b[^>]*\bsrc=[\"']([^\"']+)[\"'][^>]*>",
            flags=re.IGNORECASE,
        )

        def replace_html_image(match: re.Match) -> str:
            image_index = copy_image(match.group(1))
            return f"\n【图片{image_index}】\n" if image_index else "\n"

        markdown_text = html_image_pattern.sub(replace_html_image, markdown_text)
        markdown_text = re.sub(r"\n{4,}", "\n\n\n", markdown_text).strip()
        return markdown_text, image_urls, image_names

    def _get_content_with_mineru(self, file_path: str):
        if Path(file_path).suffix.lower() != ".pptx":
            raise ValueError("MinerU 仅支持 .pptx，不支持旧版 .ppt。")

        mineru_exe = self._resolve_mineru_executable()
        print(f"[MinerU] 开始解析 PPTX：{file_path}")

        with tempfile.TemporaryDirectory(prefix="mineru_ppt_") as temp_output_dir:
            command = [
                mineru_exe,
                "-p",
                os.path.abspath(file_path),
                "-o",
                temp_output_dir,
            ]
            try:
                result = subprocess.run(
                    command,
                    check=True,
                    capture_output=True,
                    text=True,
                    encoding="utf-8",
                    errors="replace",
                    timeout=self.mineru_timeout,
                )
            except subprocess.TimeoutExpired as error:
                raise RuntimeError(
                    f"MinerU 解析 PPTX 超时，超过 {self.mineru_timeout} 秒。"
                ) from error
            except subprocess.CalledProcessError as error:
                detail = (error.stderr or error.stdout or str(error)).strip()
                raise RuntimeError(f"MinerU 解析 PPTX 失败：{detail}") from error

            if result.stdout:
                print("[MinerU] 输出：")
                print(result.stdout[-3000:])

            try:
                markdown_file = self._find_mineru_markdown(temp_output_dir)
            except FileNotFoundError as error:
                detail = (result.stderr or result.stdout or "").strip()
                if detail:
                    detail = f"\nMinerU 输出：{detail[-3000:]}"
                raise RuntimeError(f"{error}{detail}") from error

            markdown_text = markdown_file.read_text(encoding="utf-8")
            markdown_text, image_urls, image_names = self._copy_mineru_images(
                markdown_text,
                markdown_file.parent,
            )
            markdown_text = self.ppt_noise_filter.sanitize_text(markdown_text)
            if not markdown_text:
                raise RuntimeError("MinerU 解析 PPTX 完成，但提取到的文本为空。")

            print(
                f"[MinerU] PPTX 解析成功，文本长度：{len(markdown_text)}，"
                f"图片数量：{len(image_urls)}"
            )
            return (
                markdown_text,
                image_urls,
                image_names,
                self._empty_section_image_indexes(),
            )

    def _resolve_soffice_executable(self) -> str:
        if self.ppt_soffice_exe:
            candidate = os.path.abspath(os.path.expanduser(self.ppt_soffice_exe))
            if os.path.exists(candidate):
                return candidate
        for command in ("soffice", "libreoffice"):
            resolved = shutil.which(command)
            if resolved:
                return resolved
        return ""

    def _convert_ppt_to_pdf(self, file_path: str) -> str:
        soffice = self._resolve_soffice_executable()
        if not soffice:
            raise RuntimeError(
                "未找到 LibreOffice/soffice，无法将 PPT 转 PDF 后交给 MinerU。"
                "请安装 LibreOffice，或设置 PPT_SOFFICE_EXE 指向 soffice.exe。"
            )

        output_parent = Path(self.document_base_dir) / "runtime" / "ppt_pdf"
        output_parent.mkdir(parents=True, exist_ok=True)
        if self.ppt_keep_converted_pdf:
            output_dir = output_parent / f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:8]}"
            output_dir.mkdir(parents=True, exist_ok=True)
            cleanup_context = None
        else:
            cleanup_context = tempfile.TemporaryDirectory(prefix="ppt_pdf_", dir=str(output_parent))
            output_dir = Path(cleanup_context.name)

        command = [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            os.path.abspath(file_path),
        ]
        completed = subprocess.run(
            command,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=300,
            check=False,
        )
        if completed.returncode != 0:
            if cleanup_context is not None:
                cleanup_context.cleanup()
            raise RuntimeError(
                "PPT 转 PDF 失败：{stderr}".format(stderr=completed.stderr.strip() or completed.stdout.strip())
            )

        pdf_candidates = sorted(output_dir.glob("*.pdf"), key=lambda path: path.stat().st_mtime, reverse=True)
        if not pdf_candidates:
            if cleanup_context is not None:
                cleanup_context.cleanup()
            raise RuntimeError("PPT 转 PDF 未生成 PDF 文件。")

        pdf_path = str(pdf_candidates[0])
        if cleanup_context is not None:
            # MinerU 需要在该函数返回后读取 PDF，因此先复制到持久目录。
            keep_dir = output_parent / f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:8]}"
            keep_dir.mkdir(parents=True, exist_ok=True)
            keep_pdf_path = keep_dir / pdf_candidates[0].name
            shutil.copy2(pdf_path, keep_pdf_path)
            cleanup_context.cleanup()
            pdf_path = str(keep_pdf_path)
        return pdf_path

    def _parse_with_mineru_if_available(self, file_path: str):
        if not self.ppt_to_pdf_enabled:
            print("[PPTParser] PPT_TO_PDF_ENABLED 未启用，改用原生 PPT 解析。")
            return None

        if Path(file_path).suffix.lower() == ".pptx":
            try:
                print("[PPTParser] PPTX 直接交给 MinerU 解析。")
                text, image_urls, image_names, section_image_indexes = (
                    self._get_content_with_mineru(file_path)
                )
                return self.file2document(
                    text,
                    image_urls,
                    image_names,
                    section_image_indexes,
                )
            except Exception as error:
                print(f"[PPTParser] 直接 MinerU 解析失败，改用原生 PPT 解析：{error}")
                return None

        pdf_path = None
        try:
            pdf_path = self._convert_ppt_to_pdf(file_path)
            print(f"[PPTParser] PPT 已转换为 PDF，交给 MinerU：{pdf_path}")
            from utils.PdfParser import pdf_parser

            return pdf_parser.parse_with_mineru(pdf_path, include_page_images=self.ppt_include_page_images)
        except Exception as error:
            print(f"[PPTParser] MinerU 路由失败，改用原生 PPT 解析：{error}")
            return None
        finally:
            if pdf_path and not self.ppt_keep_converted_pdf:
                shutil.rmtree(str(Path(pdf_path).parent), ignore_errors=True)

    def get_content(self, file_path):
        if not os.path.exists(file_path):
            raise FileNotFoundError(file_path)
        prs = Presentation(file_path)
        image_urls = []
        image_names = []
        layout_items = []
        base_url = os.path.join(self.document_base_dir, self.image_dir)
        slide_width = int(prs.slide_width)
        slide_height = int(prs.slide_height)
        picture_hash_counts = self._collect_picture_hash_counts(prs)

        for slide_index, slide in enumerate(prs.slides):
            slide_analysis = self._analyze_slide(
                slide,
                slide_index + 1,
                slide_width,
                slide_height,
            )
            for shape in self._iter_shapes(slide.shapes):
                try:
                    left = int(getattr(shape, "left", 0) or 0)
                    top = int(getattr(shape, "top", 0) or 0)
                    shape_text = self._shape_text(shape)
                    if shape_text:
                        item_type = "table" if getattr(shape, "has_table", False) else "text"
                        layout_items.append({
                            "type": item_type,
                            "slide": slide_index,
                            "top": top,
                            "left": left,
                            "content": shape_text,
                        })
                    if self._is_picture_shape(shape):
                        if not self._should_keep_picture_shape(
                            shape,
                            slide_analysis,
                            slide_width,
                            slide_height,
                            picture_hash_counts,
                        ):
                            continue
                        image = shape.image
                        ext = image.ext or "png"
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                        unique_filename = f"{timestamp}_{uuid.uuid4().hex}.{ext}"
                        image_path = os.path.join(base_url, unique_filename)
                        with open(image_path, "wb") as img_file:
                            img_file.write(image.blob)
                        print(f"已保存 {unique_filename}")
                        image_urls.append(image_path)
                        image_names.append(unique_filename)
                        layout_items.append({
                            "type": "image",
                            "slide": slide_index,
                            "top": top,
                            "left": left,
                            "image_index": len(image_urls),
                        })
                except Exception:
                    continue

        layout_items.sort(key=lambda item: (item["slide"], item["top"], item["left"]))
        current_section = None
        section_image_indexes = self._empty_section_image_indexes()
        text_parts = []
        current_slide = None
        for item in layout_items:
            if item["slide"] != current_slide:
                current_slide = item["slide"]
                current_section = None
                text_parts.append(f"\n【第{current_slide + 1}页】")
            if item["type"] in {"text", "table"}:
                detected_section = self._detect_section(item["content"])
                if detected_section:
                    current_section = detected_section
                text_parts.append(item["content"])
            else:
                image_index = item["image_index"]
                text_parts.append(f"【图片{image_index}】")
                if current_section:
                    section_image_indexes[current_section].append(image_index)

        return "\n".join(text_parts), image_urls, image_names, section_image_indexes

    def image_to_base64(self, image: str):
        with open(image, "rb") as f:
            image_base64 = base64.b64encode(f.read()).decode("utf-8")
            return image_base64

    def compress_image(self, image_path: str, max_size=512, pad_color=(0, 0, 0)):
        if not os.path.exists(image_path):
            raise FileNotFoundError()
        dir_name, filename = os.path.split(image_path)
        name, ext = os.path.splitext(filename)
        new_path = f"{name}_compressed_{max_size}{ext}"
        new_path = os.path.join(dir_name, new_path)

        if os.path.exists(new_path):
            return new_path

        image = Image.open(image_path).convert("RGB")
        # new_size = (448, 448)
        max_length = max(image.width, image.height)
        rate = max_size / max_length
        new_size = (int(image.width * rate), int(image.height * rate))
        resized_image = image.resize(new_size)

        new_image = Image.new("RGB", (max_size, max_size), pad_color)

        x = (max_size - new_size[0]) // 2
        y = (max_size - new_size[1]) // 2

        new_image.paste(resized_image, (x, y))

        new_image.save(new_path)
        return new_path

    def _remove_image_files(self, image_path: str):
        if os.path.exists(image_path):
            os.remove(image_path)
        dir_name, filename = os.path.split(image_path)
        name, ext = os.path.splitext(filename)
        compressed_path = os.path.join(
            dir_name,
            f"{name}_compressed_{self.model_image_max_size}{ext}",
        )
        if os.path.exists(compressed_path):
            os.remove(compressed_path)

    def generate_message(self, text, image_urls):
        text = self.ppt_noise_filter.sanitize_text(text)
        messages = []
        data = {}
        # print(text)
        prompt = """你是一位专业的设备维修分析专家。我将给你一段关于设备维修的文本（包含文字描述）以及若干张相关图片，每张图片都有唯一的编号（从1开始）且只属于一个字段。你的任务是基于这些内容，严格按照以下模板生成JSON格式的总结。请确保所有信息均来源于提供的文本和图片，不得杜撰。对于缺失的信息，对应字段留空，但不可缺失字段（文本为空字符串，图片为空列表）。
模板：
标题：<简洁的标题，点明案例名称，必须填写>
问题简介：<可包含定义解释，现象介绍，问题发生频率，后果等内容>
原因：<造成该问题的主要原因，尽量从高频到低频排序>
评估：<评估问题的手段，方法，工具等信息>
检查：<描述维修现场如何进行定位确认，要求详细>
解决方法：<现场的解决措施及根本的解决方案，要求详细>
总结：<总结问题的主要原因，后果及解决方案的关键信息>
相关图片：<与字段相关的图像，每张图像最多出现在一个字段中>

请严格按照如下JSON格式输出：
{{
    "title": "标题", // 案例名称
    "problem_intro": "问题简介文本",
    "image_urls_problem_intro": [1, 3], // 与问题简介相关的图片编号
    "causes": "原因文本",
    "image_urls_causes": [2], // 与原因相关的图片编号
    "evaluation": "评估方法文本",
    "image_urls_evaluation": [4], // 与评估相关的图片编号
    "inspection": "检查方法文本",
    "image_urls_inspection": [5], // 与检查相关的图片编号
    "solutions": "解决方案文本",
    "image_urls_solutions": [6], // 与解决方案相关的图片编号
    "key_points": "总结文本",
    "image_urls_key_points": [7, 8] // 与总结相关的图片编号
}}

注意：
1.所有内容必须严格来源于提供的文本和图片。禁止任何形式的脑补、推理、常识补充或添加原文未提及的修饰词与解释，并且最大限度利用文本。
2. 内容中不包含的信息，对应字段可以为空，若不包含图片，图片为空列表[]。
3. 内容必须基于我提供的文本和图片，图片编号从1开始，不可杜撰任何信息。
4. 你给出的回答仅包含我要求的JSON格式答案。
5. 给定图片中可能包含无关图片，请勿放进回答中。
6. 每张图片最多出现在一个字段中。
7. 内容需连贯详细，最大限度使用给定内容，请勿过分精简。
8. 各字段内容请勿大量重复，无关图片不要放入回答。

现在请分析下面的内容：
[文本内容]
{text}

[图片内容由base64给出]""".format(text=text)
        prompt += "\n注意：图片归属必须优先依据其在原文中的位置。图片位于哪个小节下，就归入对应 image_urls 字段，不得仅凭图片内容语义移动到其他字段。"
        msg_content = [{"type": "text", "text": prompt}]
        # encoding = tiktoken.get_encoding("cl100k_base")
        token_cnt = get_token_count(prompt)

        print(f"token1: {token_cnt}")
        for image in image_urls:
            print(image)
            if token_cnt >= self.input_token - 1000:
                break
            compress_image = self.compress_image(image, max_size=self.model_image_max_size)
            mime_type, _ = mimetypes.guess_type(compress_image)
            if mime_type is None:
                ext = os.path.splitext(compress_image)[1].lower()
                mime_type = {
                    '.png': 'image/png',
                    '.jpg': 'image/jpeg',
                    '.jpeg': 'image/jpeg',
                    '.webp': 'image/webp',
                    '.bmp': 'image/bmp'
                }.get(ext, 'image/jpeg')
            image_base64 = self.image_to_base64(compress_image)
            msg_content.append({
                "type": "image_url",
                "image_url": {"url": f"data:{mime_type};base64,{image_base64}"}
            })
            token_cnt += 258
            # l += len(image_base64)
        data["role"] = "user"
        data["content"] = msg_content
        messages.append(data)
        # print("len = {}".format(l))
        return messages

    def _image_indexes_to_urls(self, image_indexes, image_names):
        urls = []
        for image_index in image_indexes:
            if image_index <= 0 or image_index > len(image_names):
                continue
            urls.append(self.image_dir + "/" + image_names[image_index - 1])
        return ", ".join(urls) if urls else None

    def _apply_section_image_urls(self, result, section_image_indexes, image_names):
        if not section_image_indexes or not any(section_image_indexes.values()):
            return result, set()
        used_indexes = set()
        for section, image_indexes in section_image_indexes.items():
            result[f"image_urls_{section}"] = self._image_indexes_to_urls(image_indexes, image_names)
            used_indexes.update(image_indexes)
        return result, used_indexes
    def file2document(self, text, image_urls, image_names, section_image_indexes=None) -> Document:
        try:
            client = OpenAI(
                base_url=get_ai_base_url_alt(),
                api_key=self.api_key
            )

            messages = self.generate_message(text, image_urls)
            # print(messages)
            response = client.chat.completions.create(
                model=self.model,
                messages=messages,
                max_tokens=self.max_token
            )
            # print(response)
            ans = response.choices[0].message.content
            print(ans)
            result = json.loads(ans)
            result["title"] = normalize_document_title(result.get("title"))
            result, used_image_indexes = self._apply_section_image_urls(result, section_image_indexes, image_names)

            if not used_image_indexes:
                flag = [0] * len(image_urls)
                for key in result.keys():
                    if "image" in key:
                        image_url_content = ""
                        for image_index in result[key]:
                            if image_index > len(image_urls) or flag[image_index - 1] == 1:
                                continue
                            url = image_names[image_index - 1]
                            url = self.image_dir + "/" + url
                            image_url_content += url + ", "
                            flag[image_index - 1] = 1
                        image_url_content = image_url_content.rstrip(", ")
                        if len(result[key]) == 0:
                            image_url_content = None
                        result[key] = image_url_content
                used_image_indexes = {i + 1 for i, used in enumerate(flag) if used == 1}
            document = Document(**result,
                                is_vectorized=0)

            if len(image_urls) > 0:
                for i in range(len(image_urls)):
                    if i + 1 not in used_image_indexes:
                        self._remove_image_files(image_urls[i])

            return document

        except Exception as e:
            print(e)
            if self._is_ai_service_unavailable_error(e):
                self._set_last_error(BizCode.AI_SERVICE_UNAVAILABLE, "AI服务不可用，请稍后重试")
            else:
                self._set_last_error(BizCode.DOC_PARSE_FAILED, str(e))
            for image in image_urls:
                self._remove_image_files(image)
            return None

ppt_parser = PPTParser()

if __name__ == "__main__":
    pass
