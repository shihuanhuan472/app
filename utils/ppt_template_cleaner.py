from __future__ import annotations

from pathlib import Path
from typing import Any, Dict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn


TEMPLATE_PLACEHOLDER_TYPES_TO_REMOVE = {
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.SLIDE_NUMBER,
}


def _placeholder_type(shape) -> Any:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        return shape.placeholder_format.type
    except Exception:
        return None


def _should_remove_template_shape(shape) -> bool:
    placeholder_type = _placeholder_type(shape)
    if placeholder_type is None:
        return True
    return placeholder_type in TEMPLATE_PLACEHOLDER_TYPES_TO_REMOVE


def _should_remove_slide_placeholder(shape) -> bool:
    return _placeholder_type(shape) in TEMPLATE_PLACEHOLDER_TYPES_TO_REMOVE


def _remove_shape(shape) -> bool:
    element = getattr(shape, "_element", None)
    if element is None:
        return False
    parent = element.getparent()
    if parent is None:
        return False
    parent.remove(element)
    return True


def _remove_background(slide_like) -> bool:
    element = getattr(slide_like, "_element", None)
    if element is None:
        return False
    c_sld = element.find(qn("p:cSld"))
    if c_sld is None:
        return False
    background = c_sld.find(qn("p:bg"))
    if background is None:
        return False
    c_sld.remove(background)
    return True


def _shape_text(shape) -> str:
    parts = []
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _shape_area_ratio(shape, slide_width: int, slide_height: int) -> float:
    slide_area = max(int(slide_width) * int(slide_height), 1)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return min(max(width * height, 0) / slide_area, 1.0)


def _is_picture_shape(shape) -> bool:
    return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE


def _remove_slide_noise_shape(shape, noise_filter, slide_width: int, slide_height: int) -> tuple[bool, str]:
    text = _shape_text(shape)
    if text:
        top = int(getattr(shape, "top", 0) or 0)
        if noise_filter.is_noise_text(text, top=top, slide_height=slide_height):
            return _remove_shape(shape), "text"

    if _is_picture_shape(shape):
        try:
            area_ratio = _shape_area_ratio(shape, slide_width, slide_height)
            should_filter, reason = noise_filter.should_filter_image_bytes(
                shape.image.blob,
                area_ratio=area_ratio,
                source=f"ppt-template-cleaner:{getattr(shape, 'name', 'picture')}",
            )
        except Exception:
            should_filter, reason = False, ""
        if should_filter:
            return _remove_shape(shape), reason or "picture"

    return False, ""


def clean_pptx_template(
    source_path: str,
    output_path: str,
    *,
    remove_slide_placeholder_noise: bool = True,
    remove_slide_noise: bool = True,
    noise_filter=None,
) -> Dict[str, Any]:
    """Save a presentation copy with master/layout template noise removed."""
    source = Path(source_path)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(str(source))
    stats: Dict[str, Any] = {
        "source_path": str(source),
        "output_path": str(output),
        "masters_processed": 0,
        "layouts_processed": 0,
        "slides_processed": 0,
        "removed_template_shapes": 0,
        "removed_slide_placeholders": 0,
        "removed_backgrounds": 0,
        "removed_slide_noise_shapes": 0,
        "removed_slide_noise_text_shapes": 0,
        "removed_slide_noise_picture_shapes": 0,
    }

    for master in presentation.slide_masters:
        stats["masters_processed"] += 1
        if _remove_background(master):
            stats["removed_backgrounds"] += 1
        for shape in list(master.shapes):
            if _should_remove_template_shape(shape) and _remove_shape(shape):
                stats["removed_template_shapes"] += 1

        for layout in master.slide_layouts:
            stats["layouts_processed"] += 1
            if _remove_background(layout):
                stats["removed_backgrounds"] += 1
            for shape in list(layout.shapes):
                if _should_remove_template_shape(shape) and _remove_shape(shape):
                    stats["removed_template_shapes"] += 1

    if remove_slide_placeholder_noise:
        for slide in presentation.slides:
            stats["slides_processed"] += 1
            for shape in list(slide.shapes):
                if _should_remove_slide_placeholder(shape) and _remove_shape(shape):
                    stats["removed_slide_placeholders"] += 1

    if remove_slide_noise and noise_filter is not None and getattr(noise_filter, "enabled", False):
        slide_width = int(presentation.slide_width)
        slide_height = int(presentation.slide_height)
        for slide in presentation.slides:
            for shape in list(slide.shapes):
                removed, reason = _remove_slide_noise_shape(
                    shape,
                    noise_filter,
                    slide_width,
                    slide_height,
                )
                if not removed:
                    continue
                stats["removed_slide_noise_shapes"] += 1
                if reason == "text":
                    stats["removed_slide_noise_text_shapes"] += 1
                else:
                    stats["removed_slide_noise_picture_shapes"] += 1

    stats["changed"] = bool(
        stats["removed_template_shapes"]
        or stats["removed_slide_placeholders"]
        or stats["removed_backgrounds"]
        or stats["removed_slide_noise_shapes"]
    )
    presentation.save(str(output))
    return stats
