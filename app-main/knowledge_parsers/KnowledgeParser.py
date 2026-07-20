import csv
import os
import re
import shutil
import uuid
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from PIL import Image

from knowledge_parsers.enterprise_word_chunker import EnterpriseWordChunker
from utils.error_codes import BizCode

try:
    import pymupdf
except Exception:  # pragma: no cover - optional import is validated at runtime
    pymupdf = None


@dataclass
class KnowledgeSectionData:
    section_index: int
    section_title: str
    section_type: str = "1"
    plain_text: str = ""
    image_urls: List[str] = field(default_factory=list)
    char_start: Optional[int] = None
    char_end: Optional[int] = None
    metadata: Dict = field(default_factory=dict)


@dataclass
class KnowledgeParsedDocument:
    title: str
    summary: str
    content: str
    image_urls: List[str]
    sections: List[KnowledgeSectionData]


class KnowledgeParser:
    """
    知识库导入解析器。

    该解析器不按故障库的 problem_intro/causes/solutions 字段拆分，
    而是尽量保留原文档标题、段落和图片位置，产出知识库文档 + 章节结构。
    """

    def __init__(self):
        self.document_base_dir = os.getenv("DOCUMENT_BASE_DIR", "D:/Pycharm/code/Maintenance_Assistance_System")
        self.image_dir = os.getenv("IMAGE_DIR", "upload/images")
        self.last_error_code = None
        self.last_error_detail = None
        os.makedirs(os.path.join(self.document_base_dir, self.image_dir), exist_ok=True)

    def parse(self, file_path: str) -> Optional[KnowledgeParsedDocument]:
        self.last_error_code = None
        self.last_error_detail = None
        try:
            ext = Path(file_path).suffix.lower()
            if ext == ".pdf":
                return self._parse_pdf(file_path)
            if ext == ".docx":
                return self._parse_docx(file_path)
            if ext in {".md", ".markdown"}:
                return self._parse_markdown(file_path)
            if ext in {".html", ".mhtml"}:
                return self._parse_html(file_path)
            if ext == ".txt":
                return self._parse_txt(file_path)
            if ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp"}:
                return self._parse_image(file_path)
            if ext in {".csv"}:
                return self._parse_csv(file_path)
            if ext in {".xlsx", ".xls", ".xlsm"}:
                return self._parse_excel(file_path)
            if ext in {".pptx", ".ppt"}:
                return self._parse_ppt(file_path)
            self._set_last_error(BizCode.DOC_PARSE_FAILED, f"知识库解析器暂不支持该文件类型：{ext}")
            return None
        except Exception as exc:
            self._set_last_error(BizCode.DOC_PARSE_FAILED, str(exc))
            return None

    def _set_last_error(self, code: int, message: str):
        self.last_error_code = int(code)
        self.last_error_detail = message

    def _relative_image_path(self, filename: str) -> str:
        image_dir = self.image_dir.rstrip("/").rstrip("\\")
        return f"{image_dir}/{filename}".replace("\\", "/")

    def _copy_image_to_upload(self, source_path: str) -> str:
        ext = Path(source_path).suffix.lower() or ".png"
        filename = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex}{ext}"
        target_path = os.path.join(self.document_base_dir, self.image_dir, filename)
        shutil.copy2(source_path, target_path)
        return self._relative_image_path(filename)

    def _save_image_blob(self, blob: bytes, ext: str = ".png") -> str:
        ext = ext if ext.startswith(".") else f".{ext}"
        if ext == ".jpeg":
            ext = ".jpg"
        filename = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex}{ext}"
        target_path = os.path.join(self.document_base_dir, self.image_dir, filename)
        with open(target_path, "wb") as f:
            f.write(blob)
        return self._relative_image_path(filename)

    def _build_document(self, file_path: str, blocks: List[Dict]) -> KnowledgeParsedDocument:
        blocks = self._filter_noise_blocks(blocks)
        sections = self._blocks_to_sections(blocks, fallback_title=Path(file_path).stem)
        sections = self._with_front_matter_sections(blocks, sections)
        content = "\n\n".join(section.plain_text for section in sections if section.plain_text).strip()
        title = self._guess_title(file_path, sections, content)
        summary = self._build_summary(content)
        image_urls = []
        for section in sections:
            for image_url in section.image_urls:
                if image_url not in image_urls:
                    image_urls.append(image_url)
        return KnowledgeParsedDocument(
            title=title,
            summary=summary,
            content=content,
            image_urls=image_urls,
            sections=sections,
        )

    def _with_front_matter_sections(self, blocks: List[Dict], sections: List[KnowledgeSectionData]) -> List[KnowledgeSectionData]:
        front_sections = self._extract_front_matter_sections(blocks)
        if not front_sections:
            return sections
        existing_roles = {self._section_role(section) for section in sections}
        front_sections = [section for section in front_sections if self._section_role(section) not in existing_roles]
        insert_at = 1 if sections and self._section_role(sections[0]) == "title" else 0
        merged = sections[:insert_at] + front_sections + sections[insert_at:]
        for index, section in enumerate(merged):
            section.section_index = index
        return merged

    def _extract_front_matter_sections(self, blocks: List[Dict]) -> List[KnowledgeSectionData]:
        lines: List[str] = []
        front_tables: List[Dict] = []
        for block in blocks:
            page = block.get("page") or 1
            if page > 6:
                break
            if block.get("type") == "table":
                table_text = (block.get("text") or "").strip()
                if table_text:
                    front_tables.append(
                        {
                            "text": table_text,
                            "page": page,
                            "bbox": block.get("bbox") or (0, 0, 0, 0),
                        }
                    )
                continue
            if block.get("type") != "text":
                continue
            for raw_line in str(block.get("text") or "").splitlines():
                line = raw_line.strip()
                if line:
                    lines.append(line)

        if not lines:
            return []

        declaration_index = self._find_front_line(lines, {"宣言", "声明"})
        revision_index = self._find_front_line(lines, {"修订记录", "修订历史", "版本记录", "变更记录"})
        directory_index = self._find_front_line(lines, {"目录", "目 录", "contents"})
        if directory_index is None:
            directory_index = len(lines)

        front_sections: List[KnowledgeSectionData] = []
        if declaration_index is not None and declaration_index < directory_index:
            end_index = revision_index if revision_index is not None and revision_index > declaration_index else directory_index
            declaration_text = self._front_text(lines, declaration_index + 1, end_index)
            if declaration_text:
                front_sections.append(
                    KnowledgeSectionData(
                        section_index=0,
                        section_title="文档前置信息",
                        section_type="preface",
                        plain_text=declaration_text,
                        char_start=0,
                        char_end=len(declaration_text),
                        metadata={"image_positions": [], "section_role": "preface", "chunk_strategy": "front_matter_section_v2"},
                    )
                )

        if revision_index is not None and revision_index < directory_index:
            revision_text = self._front_table_after_heading(blocks, front_tables, {"修订记录", "修订历史", "版本记录", "变更记录"})
            if not revision_text:
                revision_text = self._front_text(lines, revision_index + 1, directory_index)
            revision_text = self._normalize_revision_history_text(revision_text)
            if revision_text:
                front_sections.append(
                    KnowledgeSectionData(
                        section_index=len(front_sections),
                        section_title="修订记录",
                        section_type="revision_history",
                        plain_text=revision_text,
                        char_start=0,
                        char_end=len(revision_text),
                        metadata={"image_positions": [], "section_role": "revision_history", "chunk_strategy": "front_matter_section_v2"},
                    )
                )
        return front_sections

    def _section_role(self, section: KnowledgeSectionData) -> str:
        return str((section.metadata or {}).get("section_role") or section.section_type or "").strip().lower()

    def _filter_noise_blocks(self, blocks: List[Dict]) -> List[Dict]:
        """过滤页眉页脚、页码、空白和明显导航噪声。

        知识库解析会把图片占位符写入正文，页脚中的保密提示和页码如果混入
        章节表，会影响查重、存储和向量召回；这里按“行”清洗后再做重复页眉
        页脚判断，避免类似“机密...请勿外传\n9”这样的组合块入库。
        """
        text_counter: Dict[str, int] = {}
        page_count = len({block.get("page") for block in blocks if block.get("page") is not None})

        cleaned_blocks: List[Dict] = []
        for block in blocks:
            if block.get("type") != "text":
                cleaned_blocks.append(block)
                continue
            cleaned_text = self._clean_text_noise_lines(block.get("text", ""))
            if not cleaned_text:
                continue
            cleaned_block = dict(block)
            cleaned_block["text"] = cleaned_text
            cleaned_blocks.append(cleaned_block)
            normalized = self._normalize_noise_text(cleaned_text)
            if normalized:
                text_counter[normalized] = text_counter.get(normalized, 0) + 1

        filtered = []
        for block in cleaned_blocks:
            if block.get("type") != "text":
                filtered.append(block)
                continue
            text = (block.get("text") or "").strip()
            normalized = self._normalize_noise_text(text)
            if not text:
                continue
            if self._is_noise_line(text):
                continue
            if page_count >= 3 and text_counter.get(normalized, 0) >= max(3, page_count // 2):
                continue
            if block.get("region") in {"header", "footer"} and len(text) <= 120:
                continue
            if normalized.lower() in {"目录", "返回顶部", "上一页", "下一页"}:
                continue
            filtered.append(block)
        return filtered

    def _clean_text_noise_lines(self, text: str) -> str:
        lines = []
        for raw_line in str(text or "").splitlines():
            line = raw_line.strip()
            if not line or self._is_noise_line(line):
                continue
            lines.append(line)
        return "\n".join(lines).strip()

    def _is_noise_line(self, text: str) -> bool:
        normalized = self._normalize_noise_text(text)
        if not normalized:
            return True
        if re.fullmatch(r"第?\s*\d+\s*页?|[-–—]?\s*\d+\s*[-–—]?", text.strip()):
            return True
        if re.fullmatch(r"page\s*\d+(\s*/\s*\d+)?", text.strip(), flags=re.IGNORECASE):
            return True
        # MGI 手册页脚常见保密提示；作为页脚噪声处理，不进入章节正文。
        if re.search(r"机密|保密|confidential", normalized, flags=re.IGNORECASE) and re.search(
            r"仅适用于|服务工程师|授权服务提供商|请勿外传|donotdistribute|internaluse",
            normalized,
            flags=re.IGNORECASE,
        ):
            return True
        return False

    def _normalize_noise_text(self, text: str) -> str:
        return re.sub(r"\s+", "", text or "")

    def _find_front_line(self, lines: List[str], keywords: set) -> Optional[int]:
        normalized_keywords = {self._normalize_noise_text(keyword).lower() for keyword in keywords}
        for index, line in enumerate(lines):
            normalized = self._normalize_noise_text(line).lower()
            if normalized in normalized_keywords:
                return index
        for index, line in enumerate(lines):
            normalized = self._normalize_noise_text(line).lower()
            if any(keyword and keyword in normalized for keyword in normalized_keywords):
                return index
        return None

    def _front_text(self, lines: List[str], start: int, end: int) -> str:
        selected = []
        for line in lines[start:end]:
            normalized = self._normalize_noise_text(line)
            if not normalized:
                continue
            if normalized in {"目录", "目錄", "contents"}:
                break
            selected.append(line)
        return "\n".join(selected).strip()

    def _front_table_after_heading(self, blocks: List[Dict], tables: List[Dict], heading_keywords: set) -> str:
        if not tables:
            return ""
        heading_markers = {self._normalize_noise_text(keyword).lower() for keyword in heading_keywords}
        heading_pos = None
        directory_pos = None
        for block in blocks:
            if (block.get("page") or 1) > 6:
                break
            if block.get("type") != "text":
                continue
            text = str(block.get("text") or "")
            for raw_line in text.splitlines():
                normalized = self._normalize_noise_text(raw_line).lower()
                if normalized in heading_markers or any(keyword and keyword in normalized for keyword in heading_markers):
                    heading_pos = self._block_position(block)
                if normalized in {"目录", "目錄", "contents", "tableofcontents"}:
                    directory_pos = self._block_position(block)
                    break
            if heading_pos and directory_pos:
                break
        if not heading_pos:
            return ""

        candidates = []
        for table in tables:
            table_pos = self._block_position(table)
            if table_pos <= heading_pos:
                continue
            if directory_pos and table_pos >= directory_pos:
                continue
            candidates.append((table_pos, table.get("text") or ""))
        candidates.sort(key=lambda item: item[0])
        return candidates[0][1].strip() if candidates else ""

    def _block_position(self, block: Dict) -> Tuple[int, float, float]:
        bbox = block.get("bbox") or (0, 0, 0, 0)
        return (int(block.get("page") or 1), float(bbox[1] if len(bbox) > 1 else 0), float(bbox[0] if bbox else 0))

    def _is_revision_table_text(self, text: str) -> bool:
        normalized = self._normalize_noise_text(text)
        return all(keyword in normalized for keyword in ("版本", "日期")) and (
            "变更说明" in normalized or "修订" in normalized
        )

    def _table_to_markdown(self, rows: List[List]) -> str:
        cleaned_rows = []
        for row in rows or []:
            cells = [str(cell or "").replace("\n", " ").strip() for cell in row]
            if any(cells):
                cleaned_rows.append(cells)
        if not cleaned_rows:
            return ""
        max_cols = max(len(row) for row in cleaned_rows)
        normalized_rows = [row + [""] * (max_cols - len(row)) for row in cleaned_rows]

        def render_row(row: List[str]) -> str:
            return "| " + " | ".join(self._strip_directory_tail(str(cell or "").replace("|", " ")) for cell in row) + " |"

        return "\n".join(
            [render_row(normalized_rows[0]), render_row(["---"] * max_cols)]
            + [render_row(row) for row in normalized_rows[1:]]
        ).strip()

    def _extract_pdf_tables(self, page, page_number: int) -> List[Dict]:
        if not hasattr(page, "find_tables"):
            return []
        extracted: List[Dict] = []
        candidates = []
        for clip in self._pdf_table_candidate_clips(page):
            candidates.append(
                {
                    "clip": clip,
                    "vertical_strategy": "text",
                    "horizontal_strategy": "lines",
                    "name": "text_columns_in_line_grid",
                }
            )
            candidates.append(
                {
                    "clip": clip,
                    "vertical_strategy": "text",
                    "horizontal_strategy": "text",
                    "name": "text_grid_in_line_region",
                }
            )
        candidates.append(
            {
                "clip": None,
                "vertical_strategy": "lines",
                "horizontal_strategy": "lines",
                "name": "line_grid",
            }
        )
        for strategy in candidates:
            try:
                find_kwargs = {
                    "vertical_strategy": strategy["vertical_strategy"],
                    "horizontal_strategy": strategy["horizontal_strategy"],
                }
                if strategy.get("clip") is not None:
                    find_kwargs["clip"] = strategy["clip"]
                tables = page.find_tables(**find_kwargs)
            except Exception:
                continue
            for table_index, table in enumerate(getattr(tables, "tables", []) or []):
                rect = tuple(getattr(table, "bbox", None) or (0, 0, 0, 0))
                rows = self._merge_table_continuation_rows(self._clean_table_rows(table.extract()))
                if strategy["name"] == "text_columns_in_line_grid":
                    rows = self._prepend_pdf_table_header_from_clip(page, strategy.get("clip"), rect, rows)
                if not self._is_valid_table_rows(rows):
                    continue
                if any(self._bbox_overlap_ratio(rect, existing.get("bbox")) > 0.85 for existing in extracted):
                    existing = next((item for item in extracted if self._bbox_overlap_ratio(rect, item.get("bbox")) > 0.85), None)
                    existing_priority = self._table_strategy_priority(existing.get("table_strategy") if existing else "")
                    current_priority = self._table_strategy_priority(strategy["name"])
                    if existing and (
                        current_priority > existing_priority
                        or (
                            current_priority == existing_priority
                            and self._table_quality_score(rows) > self._table_quality_score(existing.get("rows") or [])
                        )
                    ):
                        existing.update(
                            {
                                "text": self._table_to_markdown(rows),
                                "bbox": rect,
                                "rows": rows,
                                "table_strategy": strategy["name"],
                            }
                        )
                    continue
                table_text = self._table_to_markdown(rows)
                if not table_text:
                    continue
                extracted.append(
                    {
                        "type": "table",
                        "text": table_text,
                        "page": page_number,
                        "bbox": rect,
                        "section_type": "table_text",
                        "table_index": len(extracted),
                        "table_strategy": strategy["name"],
                        "rows": rows,
                    }
                )
        for table in extracted:
            table.pop("rows", None)
        return extracted

    def _pdf_table_candidate_clips(self, page) -> List:
        horizontal_lines = []
        try:
            drawings = page.get_drawings()
        except Exception:
            return []
        for drawing in drawings:
            rect = drawing.get("rect")
            if not rect:
                continue
            width = abs(rect.x1 - rect.x0)
            height = abs(rect.y1 - rect.y0)
            if width < max(80.0, page.rect.width * 0.25) or height > 2.0:
                continue
            horizontal_lines.append((float(rect.y0), float(rect.x0), float(rect.x1)))
        if len(horizontal_lines) < 2:
            return []
        horizontal_lines.sort(key=lambda item: item[0])
        groups = []
        current = [horizontal_lines[0]]
        for line in horizontal_lines[1:]:
            previous_y = current[-1][0]
            if line[0] - previous_y <= 140:
                current.append(line)
            else:
                if len(current) >= 2:
                    groups.append(current)
                current = [line]
        if len(current) >= 2:
            groups.append(current)

        clips = []
        for group in groups:
            x0 = max(0.0, min(line[1] for line in group) - 6.0)
            x1 = min(float(page.rect.width), max(line[2] for line in group) + 6.0)
            y0 = max(0.0, min(line[0] for line in group) - 32.0)
            y1 = min(float(page.rect.height), max(line[0] for line in group) + 16.0)
            if x1 - x0 >= 80 and y1 - y0 >= 20:
                clips.append(pymupdf.Rect(x0, y0, x1, y1))
        return clips

    def _clean_table_rows(self, rows: List[List]) -> List[List[str]]:
        cleaned_rows: List[List[str]] = []
        for row in rows or []:
            cells = [str(cell or "").replace("\n", " ").strip() for cell in row]
            if any(cells):
                cleaned_rows.append(cells)
        if not cleaned_rows:
            return []
        max_cols = max(len(row) for row in cleaned_rows)
        return [row + [""] * (max_cols - len(row)) for row in cleaned_rows]

    def _is_valid_table_rows(self, rows: List[List[str]]) -> bool:
        if len(rows) < 2:
            return False
        max_cols = max(len(row) for row in rows)
        if max_cols < 2:
            return False
        non_empty_cells = sum(1 for row in rows for cell in row if str(cell or "").strip())
        return non_empty_cells >= 4

    def _table_strategy_priority(self, strategy_name: str) -> int:
        priorities = {
            "text_columns_in_line_grid": 3,
            "line_grid": 2,
            "text_grid_in_line_region": 1,
        }
        return priorities.get(strategy_name or "", 0)

    def _prepend_pdf_table_header_from_clip(self, page, clip, table_bbox, rows: List[List[str]]) -> List[List[str]]:
        if not clip or not table_bbox or not rows:
            return rows
        if self._looks_like_table_header(rows[0]):
            return rows
        try:
            words = page.get_text("words", clip=clip)
        except Exception:
            return rows
        tx0, ty0, tx1, ty1 = table_bbox
        header_words = []
        for word in words:
            x0, y0, x1, y1, text = word[:5]
            if y1 <= ty0 + 2 and y0 >= max(0.0, ty0 - 40):
                header_words.append((x0, y0, x1, y1, str(text or "")))
        if not header_words:
            return rows
        header_words.sort(key=lambda item: (item[1], item[0]))
        line_groups: List[List[Tuple[float, float, float, float, str]]] = []
        for word in header_words:
            if not line_groups or abs(line_groups[-1][0][1] - word[1]) > 5:
                line_groups.append([word])
            else:
                line_groups[-1].append(word)
        if not line_groups:
            return rows
        header_line = max(line_groups, key=lambda group: len(group))
        columns = self._table_column_ranges_from_rows(page, rows, table_bbox)
        if len(columns) < 2:
            return rows
        header = [""] * len(columns)
        for x0, _y0, x1, _y1, text in header_line:
            center = (x0 + x1) / 2
            for index, (col_x0, col_x1) in enumerate(columns):
                if col_x0 - 4 <= center <= col_x1 + 4:
                    header[index] = " ".join(part for part in (header[index], text) if part).strip()
                    break
        if not self._looks_like_table_header(header):
            return rows
        return [header] + rows

    def _table_column_ranges_from_rows(self, page, rows: List[List[str]], table_bbox) -> List[Tuple[float, float]]:
        try:
            words = page.get_text("words", clip=pymupdf.Rect(*table_bbox))
        except Exception:
            return []
        column_count = max(len(row) for row in rows)
        if column_count < 2:
            return []
        x_values = sorted({float(word[0]) for word in words} | {float(word[2]) for word in words})
        if len(x_values) < column_count:
            x0, _y0, x1, _y1 = table_bbox
            width = (x1 - x0) / column_count
            return [(x0 + width * idx, x0 + width * (idx + 1)) for idx in range(column_count)]
        tx0, _ty0, tx1, _ty1 = table_bbox
        step = (tx1 - tx0) / column_count
        return [(tx0 + step * idx, tx0 + step * (idx + 1)) for idx in range(column_count)]

    def _looks_like_table_header(self, row: List[str]) -> bool:
        cells = [str(cell or "").strip() for cell in row]
        non_empty = [cell for cell in cells if cell]
        if len(non_empty) < 2:
            return False
        digit_heavy = sum(1 for cell in non_empty if re.search(r"\d", cell))
        long_cells = sum(1 for cell in non_empty if len(cell) > 30)
        return digit_heavy == 0 and long_cells <= max(1, len(non_empty) // 2)

    def _merge_table_continuation_rows(self, rows: List[List[str]]) -> List[List[str]]:
        if not rows:
            return []
        merged: List[List[str]] = []
        for row in rows:
            cells = [str(cell or "").strip() for cell in row]
            if not any(cells):
                continue
            first_non_empty = next((idx for idx, cell in enumerate(cells) if cell), None)
            if merged and first_non_empty is not None and first_non_empty > 0 and not cells[0]:
                target = merged[-1]
                for idx, cell in enumerate(cells):
                    if not cell:
                        continue
                    target[idx] = " ".join(part for part in (target[idx], cell) if part).strip()
                continue
            merged.append(cells)
        return merged

    def _table_quality_score(self, rows: List[List[str]]) -> int:
        if not rows:
            return 0
        non_empty_cells = sum(1 for row in rows for cell in row if str(cell or "").strip())
        non_empty_first_col = sum(1 for row in rows if row and str(row[0] or "").strip())
        return non_empty_cells + non_empty_first_col * 2 + len(rows)

    def _bbox_overlap_ratio(self, first, second) -> float:
        if not first or not second:
            return 0.0
        try:
            ax0, ay0, ax1, ay1 = first
            bx0, by0, bx1, by1 = second
            inter_w = max(0.0, min(ax1, bx1) - max(ax0, bx0))
            inter_h = max(0.0, min(ay1, by1) - max(ay0, by0))
            inter_area = inter_w * inter_h
            first_area = max(0.0, (ax1 - ax0) * (ay1 - ay0))
            second_area = max(0.0, (bx1 - bx0) * (by1 - by0))
            base_area = min(first_area, second_area)
            return inter_area / base_area if base_area else 0.0
        except Exception:
            return 0.0

    def _normalize_revision_history_text(self, text: str) -> str:
        text = str(text or "").strip()
        if not text:
            return ""
        if "|" in text and "---" in text:
            return self._trim_revision_markdown_table(text)
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if len(lines) < 8:
            return text
        normalized = [self._normalize_noise_text(line) for line in lines]
        headers = ["版本", "日期", "变更说明", "编辑"]
        header_positions = []
        for header in headers:
            header_normalized = self._normalize_noise_text(header)
            pos = next((idx for idx, value in enumerate(normalized[:8]) if value == header_normalized), None)
            if pos is None:
                return text
            header_positions.append(pos)
        data = lines[max(header_positions) + 1:]
        rows = []
        index = 0
        while index < len(data):
            version = data[index].strip()
            if not re.fullmatch(r"[A-Za-z]\d+", version):
                index += 1
                continue
            if index + 2 >= len(data):
                break
            date = data[index + 1].strip()
            change_parts = []
            editor = ""
            index += 2
            while index < len(data):
                current = data[index].strip()
                next_value = data[index + 1].strip() if index + 1 < len(data) else ""
                if re.fullmatch(r"[A-Za-z]\d+", current):
                    break
                if self._looks_like_editor(current, next_value):
                    editor = current
                    index += 1
                    break
                change_parts.append(current)
                index += 1
            rows.append([version, date, " ".join(change_parts).strip(), editor])
        return self._table_to_markdown([headers] + rows) if rows else text

    def _trim_revision_markdown_table(self, text: str) -> str:
        lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
        if len(lines) < 3:
            return text
        rows = []
        for line in lines[2:]:
            cells = [cell.strip() for cell in line.strip("|").split("|")]
            if not cells or not re.fullmatch(r"[A-Za-z]\d+", cells[0]):
                continue
            row = cells[:4] + [""] * max(0, 4 - len(cells))
            row = [self._strip_directory_tail(cell) for cell in row[:4]]
            rows.append(row[:4])
        return self._table_to_markdown([["版本", "日期", "变更说明", "编辑"]] + rows) if rows else text

    def _revision_row_contains_directory(self, cells: List[str]) -> bool:
        text = self._normalize_noise_text(" ".join(cells))
        return any(marker in text for marker in ("目录", "一般信息", "重要消息框", "上电前准备"))

    def _strip_directory_tail(self, text: str) -> str:
        value = str(text or "")
        markers = ("目录", "I 一般信息", "I一般信息", "一般信息", "1.1 重要消息框", "重要消息框", "开箱并移至实验室", "上电前准备")
        positions = [value.find(marker) for marker in markers if value.find(marker) >= 0]
        return value[: min(positions)].strip() if positions else value.strip()

    def _looks_like_editor(self, value: str, next_value: str = "") -> bool:
        value = str(value or "").strip()
        if not value or len(value) > 8:
            return False
        if value.endswith(("内容", "方法", "标准", "操作", "步骤", "逻辑", "错误", "配图")):
            return False
        if re.search(r"[，,。；;:：\d]", value):
            return False
        return bool(re.fullmatch(r"[\u4e00-\u9fffA-Za-z\s]+", value)) and (
            not next_value or re.fullmatch(r"[A-Za-z]\d+", next_value)
        )

    def _bbox_inside(self, inner, outer, tolerance: float = 2.0) -> bool:
        try:
            ix0, iy0, ix1, iy1 = inner
            ox0, oy0, ox1, oy1 = outer
            return ix0 >= ox0 - tolerance and iy0 >= oy0 - tolerance and ix1 <= ox1 + tolerance and iy1 <= oy1 + tolerance
        except Exception:
            return False

    def _extract_directory_outline(self, blocks: List[Dict]) -> List[Dict]:
        """从文档目录页提取章节顺序。

        目录项可能带编号（如 2、2.1），也可能只有标题加点线页码
        （如“上电前准备........11”）。提取后章节表按“题目 + 目录 + 目录项”生成，
        因此正常情况下章节数 = 目录标题数 + 2。
        """
        outline: List[Dict] = []
        in_directory = False
        misses_after_items = 0
        seen = set()
        pending_marker = ""

        for block in blocks[:500]:
            if block.get("type") != "text":
                continue
            for raw_line in str(block.get("text") or "").splitlines():
                line = raw_line.rstrip()
                stripped_line = line.strip()
                normalized = self._normalize_noise_text(stripped_line)
                if not line:
                    continue
                if normalized in {"目录", "目錄", "contents", "tableofcontents"}:
                    in_directory = True
                    misses_after_items = 0
                    continue
                marker_only = re.fullmatch(
                    r"(\d+(?:[.．]\d+)*|第\s*[一二三四五六七八九十百千万0-9]+\s*[章节]|[一二三四五六七八九十]+[、.．]?)",
                    stripped_line,
                )
                if marker_only:
                    pending_marker = marker_only.group(1).replace("．", ".").strip("、. ")
                    in_directory = True
                    misses_after_items = 0
                    continue
                item = self._parse_directory_line(line)
                if item:
                    if pending_marker and not item.get("marker"):
                        item["marker"] = pending_marker
                        item["display_title"] = f"{pending_marker} {item['title']}".strip()
                        item["level"] = self._directory_level(pending_marker, item["title"])
                        item["raw_text"] = f"{pending_marker}\n{item.get('raw_text') or stripped_line}"
                    pending_marker = ""
                    in_directory = True
                    misses_after_items = 0
                    key = (item.get("marker") or f"auto-{len(outline)}", self._normalize_heading_key(item["title"]))
                    if key not in seen and item["title"] and self._normalize_noise_text(item["title"]) not in {"目录", "目錄"}:
                        seen.add(key)
                        item["outline_index"] = len(outline)
                        outline.append(item)
                    continue
                pending_marker = ""
                if in_directory and outline:
                    misses_after_items += 1
                    if misses_after_items >= 8:
                        self._assign_outline_markers(outline)
                        return outline
        self._assign_outline_markers(outline)
        return outline

    def _parse_directory_line(self, line: str) -> Optional[Dict]:
        original_line = line or ""
        leading_spaces = len(original_line) - len(original_line.lstrip(" \t"))
        text = re.sub(r"\s+", " ", original_line.strip())
        if not text or self._is_noise_line(text):
            return None

        # 目录页常见形式：2.1 开箱 ........ 11 / 第2章 安装 ........ 11
        numbered_match = re.match(
            r"^(?P<marker>\d+(?:[.．]\d+)*|第\s*[一二三四五六七八九十百千万0-9]+\s*[章节]|[一二三四五六七八九十]+[、.．])\s*"
            r"(?P<title>.+?)\s*(?:[.·•…\-–—_ ]{2,}|\s+)\s*(?P<page>\d{1,4})$",
            text,
        )
        marker = ""
        title = ""
        page = None
        if numbered_match:
            marker = numbered_match.group("marker").strip().replace("．", ".")
            marker = re.sub(r"\s+", "", marker).strip("、. ")
            title = numbered_match.group("title")
            page = int(numbered_match.group("page"))
        else:
            # 无编号目录项：上电前准备................................................................11
            plain_match = re.match(r"^(?P<title>.+?)\s*(?:[.·•…\-–—_ ]{2,}|\s{2,})\s*(?P<page>\d{1,4})$", text)
            if not plain_match:
                return None
            title = plain_match.group("title")
            page = int(plain_match.group("page"))

        title = re.sub(r"[.·•…\-–—_]+", " ", title).strip(" ：:.-—–、")
        if not title or len(title) > 120:
            return None
        if self._is_noise_line(title):
            return None
        level = self._directory_level(marker, title, leading_spaces)
        display_title = f"{marker} {title}".strip() if marker else title
        return {
            "marker": marker,
            "title": title,
            "display_title": display_title,
            "level": level,
            "page": page,
            "raw_text": line.strip(),
        }

    def _directory_level(self, marker: str, title: str, leading_spaces: int = 0) -> int:
        marker = (marker or "").replace("．", ".")
        number_match = re.search(r"\d+(?:\.\d+)*", marker)
        if number_match:
            return max(1, min(6, len(number_match.group(0).split("."))))
        if re.match(r"^[一二三四五六七八九十]+[、.．]", marker or ""):
            return 1
        if leading_spaces > 0:
            return max(1, min(6, leading_spaces // 2 + 1))
        return 1

    def _assign_outline_markers(self, outline: List[Dict]):
        """为无编号目录项补齐层级编号，避免 section_type 退化成顺序号。"""
        counters = [0, 0, 0, 0, 0, 0]
        for index, item in enumerate(outline):
            marker = str(item.get("marker") or "").replace("．", ".").strip("、. ")
            numeric_match = re.search(r"\d+(?:\.\d+)*", marker)
            if numeric_match:
                numeric_marker = numeric_match.group(0)
                parts = [int(part) for part in numeric_marker.split(".") if part.isdigit()]
                for idx, part in enumerate(parts[:6]):
                    counters[idx] = part
                for idx in range(len(parts), len(counters)):
                    counters[idx] = 0
                item["computed_marker"] = numeric_marker
                item["level"] = max(1, min(6, len(parts)))
                continue

            level = max(1, min(6, int(item.get("level") or 1)))
            if level > 1 and counters[0] == 0:
                counters[0] = 1
            counters[level - 1] += 1
            for idx in range(level, len(counters)):
                counters[idx] = 0
            item["computed_marker"] = ".".join(str(part) for part in counters[:level] if part > 0) or str(index + 1)
            item["level"] = level

    def _is_directory_line(self, text: str) -> bool:
        normalized = self._normalize_noise_text(text)
        if normalized in {"目录", "目錄", "contents", "tableofcontents"}:
            return True
        lines = [line.strip() for line in str(text or "").splitlines() if line.strip()]
        return bool(lines) and all(self._parse_directory_line(line) for line in lines)

    def _normalize_heading_key(self, text: str) -> str:
        text = re.sub(r"^#{1,6}\s*", "", str(text or "").strip())
        text = re.sub(r"^\d+(?:[.．]\d+)*[、.．\s]*", "", text)
        text = re.sub(r"^第\s*[一二三四五六七八九十百千万0-9]+\s*[章节][、.．\s]*", "", text)
        text = re.sub(r"^[一二三四五六七八九十]+[、.．]\s*", "", text)
        text = re.sub(r"[\s:：.。\-–—_]+", "", text)
        return text.lower()

    def _match_directory_heading(self, text: str, outline: List[Dict], start_index: int = 0) -> Optional[int]:
        if not text or self._is_directory_line(text):
            return None
        line = str(text).splitlines()[0].strip()
        line_key = self._normalize_heading_key(line)
        line_marker = self._section_marker_from_text(line)
        search_window = outline[start_index:start_index + 30] if start_index < len(outline) else outline
        base_index = start_index if start_index < len(outline) else 0
        for offset, item in enumerate(search_window):
            title_key = self._normalize_heading_key(item.get("title", ""))
            display_key = self._normalize_heading_key(item.get("display_title", ""))
            marker = str(item.get("marker", "")).replace("．", ".").strip("、. ")
            marker_number = re.search(r"\d+(?:\.\d+)*", marker)
            marker_number = marker_number.group(0) if marker_number else marker
            marker_ok = bool(marker_number and (line.startswith(marker_number) or line_marker == marker_number))
            title_ok = bool(title_key and line_key == title_key)
            display_ok = bool(display_key and line_key == display_key)
            if title_ok or display_ok or (marker_ok and title_key and title_key in line_key):
                return base_index + offset
        return None

    def _guess_title_from_blocks(self, blocks: List[Dict], fallback_title: str) -> str:
        for block in blocks[:80]:
            if block.get("type") != "text":
                continue
            text = self._clean_text_noise_lines(block.get("text", ""))
            if not text or self._is_directory_line(text):
                continue
            first_line = text.splitlines()[0].strip().strip("#").strip()
            if first_line and len(first_line) <= 120:
                return first_line[:255]
        return (fallback_title or "未命名文档")[:255]

    def _append_text_to_section(self, section: KnowledgeSectionData, text: str, full_offset: int) -> int:
        if section.plain_text:
            section.plain_text += "\n"
            full_offset += 1
        section.plain_text += text
        return full_offset + len(text)

    def _make_outline_section(self, item: Dict, index: int, full_offset: int) -> KnowledgeSectionData:
        marker = str(item.get("marker") or "").replace("．", ".").strip("、. ")
        generated_marker = item.get("computed_marker") or marker or str(int(item.get("outline_index", max(index - 2, 0))) + 1)
        return KnowledgeSectionData(
            section_index=index,
            section_title=(item.get("title") or item.get("display_title") or "未命名章节").strip()[:255],
            section_type=generated_marker,
            plain_text="",
            char_start=full_offset,
            char_end=full_offset,
            metadata={
                "image_positions": [],
                "directory_marker": marker,
                "computed_marker": generated_marker,
                "directory_level": item.get("level", 1),
                "directory_page": item.get("page"),
                "directory_raw_text": item.get("raw_text"),
                "chunk_strategy": "directory_section_v3",
                "section_role": "body",
            },
        )

    def _blocks_to_sections_by_outline(self, blocks: List[Dict], fallback_title: str, outline: List[Dict]) -> List[KnowledgeSectionData]:
        title_text = self._guess_title_from_blocks(blocks, fallback_title)
        full_offset = 0
        title_section = KnowledgeSectionData(
            section_index=0,
            section_title=title_text,
            section_type="title",
            plain_text=title_text,
            char_start=0,
            char_end=len(title_text),
            metadata={"image_positions": [], "section_role": "title", "chunk_strategy": "directory_section_v3"},
        )
        full_offset = title_section.char_end or 0

        directory_text = "\n".join((item.get("display_title") or item.get("title") or "").strip() for item in outline if (item.get("display_title") or item.get("title")))
        if full_offset:
            full_offset += 1
        directory_section = KnowledgeSectionData(
            section_index=1,
            section_title="目录",
            section_type="directory",
            plain_text=directory_text,
            char_start=full_offset,
            char_end=full_offset + len(directory_text),
            metadata={
                "image_positions": [],
                "section_role": "directory",
                "chunk_strategy": "directory_section_v3",
                "directory_items": [
                    {
                        "marker": item.get("marker"),
                        "title": item.get("title"),
                        "display_title": item.get("display_title"),
                        "level": item.get("level"),
                        "page": item.get("page"),
                    }
                    for item in outline
                ],
            },
        )
        full_offset = directory_section.char_end or full_offset

        body_sections = [self._make_outline_section(item, index + 2, full_offset) for index, item in enumerate(outline)]
        sections = [title_section, directory_section] + body_sections
        current: Optional[KnowledgeSectionData] = None
        paragraph_index_by_section: Dict[int, int] = {}
        previous_text_by_section: Dict[int, str] = {}
        next_outline_index = 0

        for block in blocks:
            if block.get("type") == "text":
                text = (block.get("text") or "").strip()
                if not text or self._is_directory_line(text):
                    continue
                matched_index = self._match_directory_heading(text, outline, next_outline_index)
                if matched_index is not None:
                    current = body_sections[matched_index]
                    next_outline_index = max(next_outline_index, matched_index + 1)
                    continue
                if current is None:
                    continue
                if current.char_start is None:
                    current.char_start = full_offset
                full_offset = self._append_text_to_section(current, text, full_offset)
                current.char_end = full_offset
                previous_text_by_section[current.section_index] = text
                paragraph_index_by_section[current.section_index] = paragraph_index_by_section.get(current.section_index, 0) + 1
            elif block.get("type") == "image":
                if current is None:
                    continue
                image_url = block.get("image_url")
                if not image_url:
                    continue
                current.image_urls.append(image_url)
                marker = f"【图片{len(current.image_urls)}】"
                if current.char_start is None:
                    current.char_start = full_offset
                full_offset = self._append_text_to_section(current, marker, full_offset)
                current.char_end = full_offset
                paragraph_index = paragraph_index_by_section.get(current.section_index, 0)
                previous_text = previous_text_by_section.get(current.section_index, "")
                current.metadata.setdefault("image_positions", []).append(
                    {
                        "image_url": image_url,
                        "paragraph_index": max(paragraph_index - 1, 0),
                        "char_offset": full_offset,
                        "page": block.get("page"),
                        "nearby_text_before": previous_text[-200:],
                        "nearby_text_after": "",
                    }
                )
            elif block.get("type") == "table":
                if current is None:
                    continue
                text = (block.get("text") or "").strip()
                if not text:
                    continue
                if current.char_start is None:
                    current.char_start = full_offset
                full_offset = self._append_text_to_section(current, text, full_offset)
                current.char_end = full_offset
                previous_text_by_section[current.section_index] = text
                paragraph_index_by_section[current.section_index] = paragraph_index_by_section.get(current.section_index, 0) + 1

        for section in sections:
            if section.char_start is None:
                section.char_start = full_offset
            if section.char_end is None:
                section.char_end = section.char_start
        self._fill_image_after_context(sections)
        return sections

    def _blocks_to_sections(self, blocks: List[Dict], fallback_title: str) -> List[KnowledgeSectionData]:
        outline = self._extract_directory_outline(blocks)
        if outline:
            sections_by_outline = self._blocks_to_sections_by_outline(blocks, fallback_title, outline)
            if sections_by_outline:
                return sections_by_outline

        sections: List[KnowledgeSectionData] = []
        current = None
        full_offset = 0

        def start_section(title: str, section_type: str = "level_1"):
            nonlocal current
            if current and (current.plain_text.strip() or current.image_urls):
                current.char_end = full_offset
                sections.append(current)
            current = KnowledgeSectionData(
                section_index=len(sections),
                section_title=(title or fallback_title or "未命名章节").strip()[:255],
                # section_type 字段用于前端展示章节编号，如 1 / 2 / 1.1。
                # 例如“2．下阶段工作内容”会被标记为 2。
                section_type=self._normalize_section_marker(section_type, title),
                plain_text="",
                char_start=full_offset,
                metadata={"image_positions": []},
            )

        start_section(fallback_title, "level_1")
        paragraph_index = 0
        previous_text = ""

        for block in blocks:
            if block.get("type") == "text":
                text = (block.get("text") or "").strip()
                if not text:
                    continue
                if outline and self._is_directory_line(text):
                    continue
                if self._is_heading(block, text) and (current.plain_text.strip() or current.image_urls):
                    start_section(text, self._section_level_from_block(block, text))
                    paragraph_index = 0
                    previous_text = ""
                    continue
                if current.plain_text:
                    current.plain_text += "\n"
                    full_offset += 1
                current.plain_text += text
                full_offset += len(text)
                previous_text = text
                paragraph_index += 1
            elif block.get("type") == "image":
                image_url = block.get("image_url")
                if not image_url:
                    continue
                current.image_urls.append(image_url)
                marker = f"【图片{len(current.image_urls)}】"
                if current.plain_text:
                    current.plain_text += "\n"
                    full_offset += 1
                current.plain_text += marker
                full_offset += len(marker)
                current.metadata.setdefault("image_positions", []).append(
                    {
                        "image_url": image_url,
                        "paragraph_index": max(paragraph_index - 1, 0),
                        "char_offset": full_offset,
                        "page": block.get("page"),
                        "nearby_text_before": previous_text[-200:],
                        "nearby_text_after": "",
                    }
                )
            elif block.get("type") == "table":
                text = (block.get("text") or "").strip()
                if not text:
                    continue
                if current.plain_text:
                    current.plain_text += "\n"
                    full_offset += 1
                current.plain_text += text
                full_offset += len(text)
                previous_text = text
                paragraph_index += 1

        if current and (current.plain_text.strip() or current.image_urls):
            current.char_end = full_offset
            sections.append(current)

        if not sections:
            sections.append(
                KnowledgeSectionData(
                    section_index=0,
                    section_title=fallback_title or "未命名章节",
                    section_type="1",
                    plain_text="",
                    char_start=0,
                    char_end=0,
                    metadata={"image_positions": []},
                )
            )

        for index, section in enumerate(sections):
            section.section_index = index
        self._fill_missing_section_markers(sections)
        self._fill_image_after_context(sections)
        return sections

    def _normalize_section_marker(self, value: str, title: str = "") -> str:
        value = str(value or "").strip().lower()
        if re.fullmatch(r"\d+(?:\.\d+)*", value):
            return value
        marker = self._section_marker_from_text(title)
        if marker:
            return marker
        if re.fullmatch(r"level_[1-6]", value):
            return value
        return ""

    def _section_level_from_block(self, block: Dict, text: str) -> str:
        raw_level = block.get("heading_level")
        if raw_level is not None:
            try:
                return f"level_{max(1, min(6, int(raw_level)))}"
            except Exception:
                pass
        return self._section_level_from_text(text)

    def _section_level_from_text(self, text: str) -> str:
        stripped = (text or "").strip()
        markdown_match = re.match(r"^(#{1,6})\s+", stripped)
        if markdown_match:
            return f"level_{len(markdown_match.group(1))}"
        number_match = re.match(r"^(\d+(?:[.．]\d+)*)[、.．\s]", stripped)
        if number_match:
            level = len(re.split(r"[.．]", number_match.group(1)))
            return f"level_{max(1, min(6, level))}"
        if re.match(r"^[一二三四五六七八九十]+[、.．\s]", stripped):
            return "level_1"
        return "level_1"

    def _section_marker_from_text(self, text: str) -> str:
        stripped = (text or "").strip()
        number_match = re.match(r"^(\d+(?:[.．]\d+)*)[、.．\s]", stripped)
        if number_match:
            return number_match.group(1).replace("．", ".").strip(".")
        chinese_match = re.match(r"^([一二三四五六七八九十]+)[、.．\s]", stripped)
        if chinese_match:
            value = self._chinese_number_to_int(chinese_match.group(1))
            if value:
                return str(value)
        return ""

    def _chinese_number_to_int(self, value: str) -> Optional[int]:
        digits = {"一": 1, "二": 2, "三": 3, "四": 4, "五": 5, "六": 6, "七": 7, "八": 8, "九": 9}
        if value == "十":
            return 10
        if value.startswith("十"):
            return 10 + digits.get(value[1:], 0)
        if "十" in value:
            left, _, right = value.partition("十")
            return digits.get(left, 0) * 10 + digits.get(right, 0)
        return digits.get(value)

    def _fill_missing_section_markers(self, sections: List[KnowledgeSectionData]):
        counters = [0, 0, 0, 0, 0, 0]
        for index, section in enumerate(sections):
            raw = str(section.section_type or "").strip().lower()
            if raw in {"title", "directory"}:
                section.section_type = raw
                continue
            if re.fullmatch(r"\d+(?:\.\d+)*", raw):
                parts = [int(part) for part in raw.split(".") if part.isdigit()]
                for idx, part in enumerate(parts[:6]):
                    counters[idx] = part
                for idx in range(len(parts), len(counters)):
                    counters[idx] = 0
                section.section_type = raw
                continue

            level_match = re.fullmatch(r"level_([1-6])", raw)
            level = int(level_match.group(1)) if level_match else 1
            counters[level - 1] += 1
            for idx in range(level, len(counters)):
                counters[idx] = 0
            if level > 1 and counters[0] == 0:
                counters[0] = 1
            section.section_type = ".".join(str(part) for part in counters[:level] if part > 0) or str(index + 1)

    def _fill_image_after_context(self, sections: List[KnowledgeSectionData]):
        for section in sections:
            paragraphs = [p.strip() for p in section.plain_text.splitlines() if p.strip()]
            for position in section.metadata.get("image_positions", []):
                paragraph_index = position.get("paragraph_index", 0)
                if paragraph_index + 1 < len(paragraphs):
                    position["nearby_text_after"] = paragraphs[paragraph_index + 1][:200]

    def _is_heading(self, block: Dict, text: str) -> bool:
        if block.get("is_heading"):
            return True
        stripped = text.strip()
        if len(stripped) > 80:
            return False
        if re.match(r"^#{1,6}\s+", stripped):
            return True
        if re.match(r"^(\d+(\.\d+)*|[一二三四五六七八九十]+)[、.．\s]", stripped):
            return True
        if stripped.endswith(("：", ":")) and len(stripped) <= 30:
            return True
        return False

    def _guess_title(self, file_path: str, sections: List[KnowledgeSectionData], content: str) -> str:
        if sections:
            first_section_text = (sections[0].plain_text or "").strip()
            if first_section_text:
                first_line = first_section_text.splitlines()[0].strip().strip("#").strip()
                if first_line and len(first_line) <= 80:
                    return first_line[:255]
        for section in sections:
            title = (section.section_title or "").strip()
            if title:
                return title[:255]
        for line in content.splitlines():
            line = line.strip().strip("#").strip()
            if line:
                return line[:255]
        return Path(file_path).stem[:255]

    def _build_summary(self, content: str) -> str:
        text = re.sub(r"\s+", " ", content or "").strip()
        return text[:500]

    def _parse_docx(self, file_path: str) -> KnowledgeParsedDocument:
        chunker = EnterpriseWordChunker(save_image_blob=self._save_image_blob)
        chunks = chunker.parse(file_path)
        sections: List[KnowledgeSectionData] = []
        full_offset = 0
        image_urls = []

        for index, chunk in enumerate(chunks):
            plain_text = chunk.plain_text or ""
            section_image_urls = list(chunk.image_urls or [])
            for image_url in section_image_urls:
                if image_url not in image_urls:
                    image_urls.append(image_url)
            section = KnowledgeSectionData(
                section_index=index,
                section_title=chunk.title or Path(file_path).stem,
                section_type=chunk.section_type or str(index + 1),
                plain_text=plain_text,
                image_urls=section_image_urls,
                char_start=full_offset,
                char_end=full_offset + len(plain_text),
                metadata=chunk.metadata or {},
            )
            sections.append(section)
            full_offset = section.char_end + 1

        content = "\n\n".join(section.plain_text for section in sections if section.plain_text).strip()
        title = self._guess_title(file_path, sections, content)
        return KnowledgeParsedDocument(
            title=title,
            summary=self._build_summary(content),
            content=content,
            image_urls=image_urls,
            sections=sections,
        )

    def _docx_heading_level(self, style_name: str) -> Optional[int]:
        style_name = style_name or ""
        match = re.search(r"(?:heading|标题)\s*([1-6])", style_name)
        if match:
            return int(match.group(1))
        if style_name.startswith(("heading", "标题")):
            return 1
        return None

    def _iter_docx_blocks(self, doc):
        for child in doc.element.body.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, doc)
            elif isinstance(child, CT_Tbl):
                yield Table(child, doc)

    def _docx_images_from_element(self, doc, element) -> List[Dict]:
        blocks = []
        rel_ids = []
        for xpath in (".//a:blip/@r:embed", ".//a:blip/@r:link"):
            try:
                rel_ids.extend(element.xpath(xpath))
            except Exception:
                continue
        for rel_id in rel_ids:
            if rel_id not in doc.part.rels:
                continue
            rel = doc.part.rels[rel_id]
            if "image" not in rel.target_ref:
                continue
            image_part = rel.target_part
            ext = image_part.content_type.split("/")[-1]
            image_url = self._save_image_blob(image_part.blob, ext)
            blocks.append({"type": "image", "image_url": image_url})
        return blocks

    def _docx_table_to_text(self, table: Table) -> str:
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
        return self._table_to_markdown(rows)

    def _parse_pdf(self, file_path: str) -> KnowledgeParsedDocument:
        if pymupdf is None:
            raise RuntimeError("pymupdf 未安装，无法解析 PDF")
        doc = pymupdf.open(file_path)
        blocks: List[Dict] = []
        for page_index in range(len(doc)):
            page = doc[page_index]
            page_height = page.rect.height or 1
            table_blocks = self._extract_pdf_tables(page, page_index + 1)
            table_bboxes = [table.get("bbox") for table in table_blocks if table.get("bbox")]
            for table_block in table_blocks:
                table_block["block_order"] = len(blocks)
                blocks.append(table_block)
            text_dict = page.get_text("dict")
            for block in text_dict.get("blocks", []):
                if block.get("type") != 0:
                    continue
                text, max_size = self._pdf_block_text_and_size(block)
                if not text:
                    continue
                x0, y0, x1, y1 = block.get("bbox", (0, 0, 0, 0))
                if any(self._bbox_inside((x0, y0, x1, y1), table_bbox) for table_bbox in table_bboxes):
                    continue
                region = "body"
                if y0 < page_height * 0.06:
                    region = "header"
                elif y1 > page_height * 0.94:
                    region = "footer"
                blocks.append(
                    {
                        "type": "text",
                        "text": text,
                        "page": page_index + 1,
                        "region": region,
                        "bbox": (x0, y0, x1, y1),
                        "block_order": len(blocks),
                        "is_heading": len(text) <= 80 and max_size >= 13,
                    }
                )
            for img in page.get_images(full=True):
                xref = img[0]
                rects = page.get_image_rects(xref)
                if not rects:
                    continue
                pix = pymupdf.Pixmap(doc, xref)
                if pix.n - pix.alpha > 3:
                    pix = pymupdf.Pixmap(pymupdf.csRGB, pix)
                image_bytes = pix.tobytes("png")
                image_url = self._save_image_blob(image_bytes, ".png")
                pix = None
                for rect in rects[:1]:
                    blocks.append(
                        {
                            "type": "image",
                            "image_url": image_url,
                            "page": page_index + 1,
                            "bbox": (rect.x0, rect.y0, rect.x1, rect.y1),
                            "block_order": len(blocks),
                        }
                    )
        blocks.sort(key=lambda item: (item.get("page") or 0, (item.get("bbox") or (0, 0, 0, 0))[1], (item.get("bbox") or (0, 0, 0, 0))[0], item.get("block_order") or 0))
        return self._build_document(file_path, blocks)

    def _pdf_block_text_and_size(self, block) -> Tuple[str, float]:
        lines = []
        max_size = 0.0
        for line in block.get("lines", []):
            spans = []
            for span in line.get("spans", []):
                span_text = span.get("text", "")
                if span_text:
                    spans.append(span_text)
                max_size = max(max_size, float(span.get("size", 0) or 0))
            line_text = "".join(spans).strip()
            if line_text:
                lines.append(line_text)
        return "\n".join(lines).strip(), max_size

    def _parse_markdown(self, file_path: str) -> KnowledgeParsedDocument:
        text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        blocks: List[Dict] = []
        base_dir = Path(file_path).parent
        image_pattern = re.compile(r"!\[[^\]]*]\(([^)]+)\)")
        for raw_line in text.splitlines():
            line = raw_line.strip()
            if not line:
                continue
            matched_images = image_pattern.findall(line)
            clean_line = image_pattern.sub("", line).strip()
            if clean_line:
                heading_match = re.match(r"^(#{1,6})\s+", line)
                blocks.append({
                    "type": "text",
                    "text": clean_line.lstrip("#").strip(),
                    "is_heading": bool(heading_match),
                    "heading_level": len(heading_match.group(1)) if heading_match else None,
                })
            for image_ref in matched_images:
                image_ref = image_ref.strip().strip('"').strip("'")
                image_path = (base_dir / image_ref).resolve()
                if image_path.exists():
                    blocks.append({"type": "image", "image_url": self._copy_image_to_upload(str(image_path))})
        return self._build_document(file_path, blocks)

    def _parse_html(self, file_path: str) -> KnowledgeParsedDocument:
        html = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
            tag.decompose()
        blocks: List[Dict] = []
        for element in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "img"]):
            if element.name == "img":
                src = element.get("src") or ""
                if src.startswith("data:"):
                    continue
                image_path = (Path(file_path).parent / src).resolve()
                if image_path.exists():
                    blocks.append({"type": "image", "image_url": self._copy_image_to_upload(str(image_path))})
                continue
            text = element.get_text(" ", strip=True)
            if text:
                heading_level = int(element.name[1]) if element.name.startswith("h") and element.name[1:].isdigit() else None
                blocks.append({"type": "text", "text": text, "is_heading": heading_level is not None, "heading_level": heading_level})
        return self._build_document(file_path, blocks)

    def _parse_txt(self, file_path: str) -> KnowledgeParsedDocument:
        text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        blocks = [{"type": "text", "text": line.strip()} for line in text.splitlines() if line.strip()]
        return self._build_document(file_path, blocks)

    def _parse_image(self, file_path: str) -> KnowledgeParsedDocument:
        image_url = self._copy_image_to_upload(file_path)
        title = Path(file_path).stem
        section = KnowledgeSectionData(
            section_index=0,
            section_title=title,
            section_type="1",
            plain_text=f"图片资料：{title}\n【图片1】",
            image_urls=[image_url],
            char_start=0,
            char_end=len(title),
            metadata={
                "image_positions": [
                    {
                        "image_url": image_url,
                        "paragraph_index": 0,
                        "char_offset": len(title),
                        "nearby_text_before": title,
                        "nearby_text_after": "",
                    }
                ]
            },
        )
        return KnowledgeParsedDocument(title=title, summary=title, content=section.plain_text, image_urls=[image_url], sections=[section])

    def _parse_csv(self, file_path: str) -> KnowledgeParsedDocument:
        rows = []
        with open(file_path, "r", encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.reader(f)
            for index, row in enumerate(reader):
                if index >= 200:
                    break
                rows.append(" | ".join(cell.strip() for cell in row if cell.strip()))
        blocks = [{"type": "text", "text": "\n".join(row for row in rows if row)}]
        return self._build_document(file_path, blocks)

    def _parse_excel(self, file_path: str) -> KnowledgeParsedDocument:
        try:
            import openpyxl
        except Exception:
            raise RuntimeError("openpyxl 未安装，无法解析 Excel")
        workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        blocks: List[Dict] = []
        for sheet in workbook.worksheets[:5]:
            blocks.append({"type": "text", "text": sheet.title, "is_heading": True})
            for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
                if row_index >= 100:
                    break
                line = " | ".join(str(cell).strip() for cell in row if cell is not None and str(cell).strip())
                if line:
                    blocks.append({"type": "text", "text": line})
        return self._build_document(file_path, blocks)

    def _parse_ppt(self, file_path: str) -> KnowledgeParsedDocument:
        try:
            from pptx import Presentation
        except Exception:
            raise RuntimeError("python-pptx 未安装，无法解析 PPT")
        presentation = Presentation(file_path)
        blocks: List[Dict] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            blocks.append({"type": "text", "text": f"第{slide_index}页", "is_heading": True})
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        blocks.append({"type": "text", "text": text})
                if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                    ext = shape.image.ext or "png"
                    image_url = self._save_image_blob(shape.image.blob, ext)
                    blocks.append({"type": "image", "image_url": image_url})
        return self._build_document(file_path, blocks)


knowledge_parser = KnowledgeParser()
