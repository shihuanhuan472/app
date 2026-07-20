import argparse
import asyncio
import contextlib
import hashlib
import io
import json
import os
import re
import shutil
import sys
import time
from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path
from types import SimpleNamespace
from typing import Any, Dict, Iterable, List, Optional, Tuple

from dotenv import load_dotenv


PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

load_dotenv(PROJECT_ROOT / ".env")
os.environ.setdefault("HF_ENDPOINT", "https://hf-mirror.com")
os.environ.setdefault("TRANSFORMERS_OFFLINE", "1")
os.environ.setdefault("HF_DATASETS_OFFLINE", "1")
os.environ.setdefault("HF_HUB_OFFLINE", "1")
model_download_url = os.getenv("MODEL_DOWNLOAD_URL")
if model_download_url:
    os.environ["HF_HOME"] = model_download_url


DEFAULT_UPLOAD_ROOT = PROJECT_ROOT / "upload" / "source_documents"
if not DEFAULT_UPLOAD_ROOT.exists():
    DEFAULT_UPLOAD_ROOT = PROJECT_ROOT / "upload"

SUPPORTED_TEXT_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".txt",
    ".md",
    ".html",
    ".htm",
    ".mhtml",
    ".csv",
    ".xlsx",
    ".xls",
}

RAG_METRIC_NAMES = [
    "context_precision",
    "context_recall",
    "faithfulness",
    "answer_accuracy",
]


def now_stamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def rel_path(path: Path) -> str:
    try:
        return str(path.resolve().relative_to(PROJECT_ROOT.resolve())).replace("\\", "/")
    except Exception:
        return str(path).replace("\\", "/")


def ensure_parent(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def load_json_list(path: Path) -> List[Dict[str, Any]]:
    if not path.exists() or path.stat().st_size == 0:
        return []
    with path.open("r", encoding="utf-8-sig") as f:
        data = json.load(f)
    if not isinstance(data, list):
        raise ValueError(f"{path} must contain a JSON list")
    return data


def write_json(path: Path, data: Any) -> None:
    ensure_parent(path)
    with path.open("w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def stable_case_id(prefix: str, *parts: Any) -> str:
    raw = "\n".join(str(part or "") for part in parts)
    digest = hashlib.sha1(raw.encode("utf-8")).hexdigest()[:12]
    return f"{prefix}-{digest}"


def normalize_extensions(raw: str) -> set:
    if not raw:
        return set(SUPPORTED_TEXT_EXTENSIONS)
    values = set()
    for item in raw.split(","):
        item = item.strip().lower()
        if not item:
            continue
        if not item.startswith("."):
            item = "." + item
        values.add(item)
    return values


def normalize_image_list(value: Any) -> List[str]:
    if not value:
        return []
    if isinstance(value, (list, tuple)):
        result: List[str] = []
        for item in value:
            result.extend(normalize_image_list(item))
        return result
    text = str(value).strip()
    if not text:
        return []
    try:
        parsed = json.loads(text)
        if isinstance(parsed, (list, tuple)):
            return normalize_image_list(parsed)
    except Exception:
        pass
    images = [item.strip().replace("\\", "/") for item in text.split(",") if item.strip()]
    deduped = []
    seen = set()
    for image in images:
        if image not in seen:
            deduped.append(image)
            seen.add(image)
    return deduped


def image_path_exists(image: str) -> bool:
    if not image:
        return False
    path = Path(image)
    if not path.is_absolute():
        path = PROJECT_ROOT / image
    return path.exists()


def resolve_project_path(path_text: str) -> Path:
    path = Path(path_text)
    if not path.is_absolute():
        path = PROJECT_ROOT / path
    return path


def existing_image_list(value: Any) -> List[str]:
    images = []
    seen = set()
    for image in normalize_image_list(value):
        if image in seen or not image_path_exists(image):
            continue
        images.append(image)
        seen.add(image)
    return images


def copy_eval_images(images: List[str], image_dir: Path, case_seed: str, max_images: int) -> List[str]:
    image_dir.mkdir(parents=True, exist_ok=True)
    copied = []
    for index, image in enumerate(images[: max(1, max_images)], start=1):
        source_path = resolve_project_path(image)
        suffix = source_path.suffix or ".jpg"
        digest = hashlib.sha1(f"{case_seed}\n{image}\n{index}".encode("utf-8")).hexdigest()[:12]
        target_name = f"{digest}_{index}{suffix.lower()}"
        target_path = image_dir / target_name
        if not target_path.exists():
            shutil.copy2(source_path, target_path)
        copied.append(rel_path(target_path))
    return copied


def select_query_images(item: Dict[str, Any], use_source_images: bool, max_query_images: int) -> str:
    images = normalize_image_list(item.get("query_images") or item.get("user_uploaded_images"))
    if not images and use_source_images:
        images = normalize_image_list(item.get("source_images"))
    images = [image for image in images if image_path_exists(image)]
    if max_query_images > 0:
        images = images[:max_query_images]
    return ", ".join(images)


def iter_files(root: Path, extensions: set) -> Iterable[Path]:
    if not root.exists():
        return []
    return (
        path
        for path in root.rglob("*")
        if path.is_file() and path.suffix.lower() in extensions
    )


def trim_text(text: str, max_chars: int) -> str:
    text = re.sub(r"\n{3,}", "\n\n", str(text or "")).strip()
    if len(text) <= max_chars:
        return text
    part = max(1, max_chars // 3)
    middle_start = max(0, (len(text) // 2) - (part // 2))
    return "\n\n".join(
        [
            text[:part].strip(),
            "[middle omitted]",
            text[middle_start : middle_start + part].strip(),
            "[tail omitted]",
            text[-part:].strip(),
        ]
    )


def read_plain_file(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gbk", "gb18030"):
        try:
            return path.read_text(encoding=encoding, errors="ignore")
        except Exception:
            continue
    return path.read_text(errors="ignore")


def extract_text_from_file(path: Path, max_chars: int) -> Tuple[str, Optional[str]]:
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            import pymupdf

            doc = pymupdf.open(str(path))
            text = "\n".join(page.get_text().strip() for page in doc)
        elif ext == ".docx":
            from docx import Document as DocxDocument

            doc = DocxDocument(str(path))
            paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
            table_cells = []
            for table in doc.tables:
                for row in table.rows:
                    values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if values:
                        table_cells.append(" | ".join(values))
            text = "\n".join(paragraphs + table_cells)
        elif ext == ".pptx":
            from pptx import Presentation

            prs = Presentation(str(path))
            parts = []
            for index, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text.strip())
                if slide_text:
                    parts.append(f"Slide {index}\n" + "\n".join(slide_text))
            text = "\n\n".join(parts)
        elif ext in {".html", ".htm", ".mhtml"}:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(read_plain_file(path), "html.parser")
            text = soup.get_text("\n")
        elif ext in {".xlsx", ".xls"}:
            import pandas as pd

            frames = pd.read_excel(path, sheet_name=None, nrows=80)
            parts = []
            for sheet_name, df in frames.items():
                parts.append(f"[sheet: {sheet_name}]")
                parts.append(df.fillna("").to_csv(index=False))
            text = "\n".join(parts)
        elif ext == ".csv":
            import pandas as pd

            df = pd.read_csv(path, nrows=200)
            text = df.fillna("").to_csv(index=False)
        elif ext in {".txt", ".md"}:
            text = read_plain_file(path)
        else:
            return "", f"unsupported extension: {ext}"
    except Exception as exc:
        return "", f"{type(exc).__name__}: {exc}"

    text = trim_text(text, max_chars=max_chars)
    if not text:
        return "", "empty text"
    return text, None


def strip_code_fence(text: str) -> str:
    text = str(text or "").strip()
    text = re.sub(r"^```(?:json)?\s*", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\s*```$", "", text)
    return text.strip()


def extract_json_array(text: str) -> List[Dict[str, Any]]:
    cleaned = strip_code_fence(text)
    start = cleaned.find("[")
    end = cleaned.rfind("]")
    if start >= 0 and end > start:
        cleaned = cleaned[start : end + 1]
    data = json.loads(cleaned)
    if not isinstance(data, list):
        raise ValueError("expected JSON array")
    return [item for item in data if isinstance(item, dict)]


def extract_json_object(text: str) -> Dict[str, Any]:
    cleaned = strip_code_fence(text)
    start = cleaned.find("{")
    end = cleaned.rfind("}")
    if start >= 0 and end > start:
        cleaned = cleaned[start : end + 1]
    data = json.loads(cleaned)
    if not isinstance(data, dict):
        raise ValueError("expected JSON object")
    return data


def resolve_ai_config(
    base_url: Optional[str],
    api_key: Optional[str],
    model: Optional[str],
    *,
    base_url_env: str = "AI_BASE_URL",
    model_env: str = "MODEL_AI",
) -> Tuple[str, str, str]:
    load_dotenv(PROJECT_ROOT / ".env")
    if not base_url:
        base_url = os.getenv(base_url_env)
    if not base_url:
        from utils.ai_endpoint import get_ai_base_url

        base_url = get_ai_base_url()
    if not api_key:
        api_key = os.getenv("API_KEY", "EMPTY")
    if not model:
        model = os.getenv(model_env, "/models/Qwen3-VL-8B-Instruct")
    return str(base_url).rstrip("/"), api_key, model


def chat_completion(
    messages: List[Dict[str, Any]],
    *,
    base_url: str,
    api_key: str,
    model: str,
    max_tokens: int,
    temperature: float = 0.0,
) -> str:
    from openai import OpenAI

    client = OpenAI(base_url=base_url, api_key=api_key)
    response = client.chat.completions.create(
        model=model,
        messages=messages,
        max_tokens=max_tokens,
        temperature=temperature,
    )
    return response.choices[0].message.content or ""


def build_dataset_prompt(title: str, source: str, text: str, questions_per_doc: int) -> List[Dict[str, Any]]:
    return [
        {
            "role": "system",
            "content": (
                "你是企业知识库问答评测集构造专家。你只根据给定文档生成评测问题和标准答案，"
                "问题必须能从文档中直接回答，标准答案必须简洁准确。"
            ),
        },
        {
            "role": "user",
            "content": (
                f"请基于下面文档生成 {questions_per_doc} 个中文问答对，用于评估RAG问答准确率。\n"
                "要求：\n"
                "1. 问题要包含具体设备、模块、故障现象、参数或操作名，避免“本文/该问题/上述”等代词。\n"
                "2. 标准答案只写文档支持的关键事实，不要扩展。\n"
                "3. 输出严格为JSON数组，不要Markdown，不要解释。\n"
                "4. 字符串内部不要使用英文双引号，如需引用请使用中文引号。\n"
                '5. 每个对象格式为：{"question":"...","ground_truth":"..."}。\n\n'
                f"文档标题：{title}\n"
                f"文档来源：{source}\n"
                f"文档内容：\n{text}"
            ),
        },
    ]


def build_json_repair_prompt(raw_response: str) -> List[Dict[str, Any]]:
    return [
        {
            "role": "system",
            "content": "你是JSON修复器。只修复格式，不改写问题和答案含义。",
        },
        {
            "role": "user",
            "content": (
                "下面内容本应是问答对JSON数组，但格式可能不合法。"
                "请修复为严格JSON数组，不要Markdown，不要解释。"
                '每个对象只保留两个字段：{"question":"...","ground_truth":"..."}。\n\n'
                f"待修复内容：\n{raw_response}"
            ),
        },
    ]


def build_judge_prompt(question: str, ground_truth: str, answer: str) -> List[Dict[str, Any]]:
    return [
        {
            "role": "system",
            "content": (
                "你是公正的企业问答评估裁判。请比较模型答案和标准答案，"
                "只依据标准答案判断核心事实是否回答正确。"
            ),
        },
        {
            "role": "user",
            "content": (
                "请输出严格JSON对象，不要Markdown。格式：\n"
                '{"score":0.0,"passed":false,"reason":"简短原因"}\n\n'
                "评分规则：\n"
                "1. score取0到1。核心事实完整且无冲突为1；部分正确为0.5左右；错误、无关或无法回答为0。\n"
                "2. passed表示score是否达到可展示准确答案水平。\n"
                "3. 若模型答案多出未被标准答案支持且可能误导的信息，应扣分。\n\n"
                f"问题：{question}\n"
                f"标准答案：{ground_truth}\n"
                f"模型答案：{answer}"
            ),
        },
    ]


def build_rag_metrics_prompt(
    question: str,
    ground_truth: str,
    answer: str,
    contexts: str,
) -> List[Dict[str, Any]]:
    return [
        {
            "role": "system",
            "content": (
                "你是企业RAG系统评估专家。请严格根据给定问题、标准答案、模型答案和检索上下文打分。"
                "所有分数范围为0到1，1表示最好，0表示最差。只输出JSON对象，不要Markdown。"
            ),
        },
        {
            "role": "user",
            "content": (
                "请输出如下JSON格式：\n"
                "{"
                '"context_precision":0.0,'
                '"context_recall":0.0,'
                '"faithfulness":0.0,'
                '"answer_accuracy":0.0,'
                '"reason":"简短说明"'
                "}\n\n"
                "指标定义：\n"
                "1. context_precision：上下文精度。检索上下文中与问题和标准事实相关的内容是否排在前面、噪声是否少。\n"
                "2. context_recall：上下文召回率。检索上下文是否覆盖标准答案所需的关键事实。\n"
                "3. faithfulness：忠实度。模型答案是否被检索上下文支持，是否编造或加入上下文没有的信息。\n"
                "4. answer_accuracy：答案准确性。模型答案相对标准答案是否正确、完整、无冲突。\n\n"
                "打分参考：\n"
                "- 1.0：完全满足；0.7：基本满足但有轻微缺失；0.5：部分满足；0.0：错误、无关或无法判断。\n"
                "- faithfulness只看答案是否被上下文支持；answer_accuracy只看答案是否回答对标准答案。\n\n"
                f"问题：{question}\n\n"
                f"标准答案：{ground_truth}\n\n"
                f"模型答案：{answer}\n\n"
                f"检索上下文：\n{contexts or '无'}"
            ),
        },
    ]


def normalize_qa_item(
    item: Dict[str, Any],
    *,
    source_prefix: str,
    source_title: str,
    source_file: str = "",
    source_doc_id: Optional[int] = None,
    source_library_type: Optional[str] = None,
    source_images: Optional[List[str]] = None,
) -> Optional[Dict[str, Any]]:
    question = str(item.get("question") or "").strip()
    ground_truth = str(item.get("ground_truth") or item.get("answer") or "").strip()
    if not question or not ground_truth:
        return None

    source_key = stable_case_id(
        source_prefix,
        source_title,
        source_file,
        source_doc_id,
        source_library_type,
    )
    qa = {
        "id": stable_case_id(source_prefix, source_title, source_file, question),
        "source_key": source_key,
        "question": question,
        "ground_truth": ground_truth,
        "source_title": source_title,
        "source_file": source_file,
    }
    images = normalize_image_list(source_images)
    if images:
        qa["source_images"] = images
    if source_doc_id is not None:
        qa["source_doc_id"] = int(source_doc_id)
    if source_library_type:
        qa["source_library_type"] = source_library_type
    return qa


def file_inventory(upload_root: Path) -> Dict[str, Any]:
    if not upload_root.exists():
        return {"upload_root": rel_path(upload_root), "exists": False}

    by_ext = Counter()
    by_top_dir = Counter()
    total_size = 0
    for path in upload_root.rglob("*"):
        if not path.is_file():
            continue
        by_ext[path.suffix.lower() or "(none)"] += 1
        total_size += path.stat().st_size
        try:
            top = path.relative_to(upload_root).parts[0]
        except Exception:
            top = "."
        by_top_dir[top] += 1

    return {
        "upload_root": rel_path(upload_root),
        "exists": True,
        "file_count": sum(by_ext.values()),
        "total_size_mb": round(total_size / 1024 / 1024, 2),
        "by_extension": dict(sorted(by_ext.items())),
        "by_top_dir": dict(sorted(by_top_dir.items())),
    }


async def db_inventory() -> Dict[str, Any]:
    from sqlalchemy import func, select

    from database import AsyncSessionLocal
    from models import DocumentBreakdown, DocumentKnowledge, SourceDocument

    async with AsyncSessionLocal() as session:
        result: Dict[str, Any] = {}
        for name, model in (
            ("breakdown", DocumentBreakdown),
            ("knowledge", DocumentKnowledge),
        ):
            total = await session.scalar(
                select(func.count()).select_from(model).where(model.is_deleted == 0)
            )
            vectorized = await session.scalar(
                select(func.count())
                .select_from(model)
                .where(model.is_deleted == 0, model.is_vectorized == 1)
            )
            result[name] = {
                "active_documents": int(total or 0),
                "vectorized_documents": int(vectorized or 0),
                "not_vectorized_documents": int((total or 0) - (vectorized or 0)),
            }

        status_rows = await session.execute(
            select(SourceDocument.status, func.count())
            .where(SourceDocument.is_deleted == 0)
            .group_by(SourceDocument.status)
        )
        result["source_documents_by_status"] = {
            str(status or "unknown"): int(count) for status, count in status_rows.all()
        }
        return result


async def collect_db_documents(limit_docs: Optional[int], only_vectorized: bool) -> List[Dict[str, Any]]:
    from sqlalchemy import select

    from database import AsyncSessionLocal
    from models import DocumentBreakdown, DocumentKnowledge, KnowledgeDocumentSection

    docs: List[Dict[str, Any]] = []
    async with AsyncSessionLocal() as session:
        breakdown_stmt = select(DocumentBreakdown).where(DocumentBreakdown.is_deleted == 0)
        if only_vectorized:
            breakdown_stmt = breakdown_stmt.where(DocumentBreakdown.is_vectorized == 1)
        breakdown_stmt = breakdown_stmt.order_by(DocumentBreakdown.id.asc())
        if limit_docs:
            breakdown_stmt = breakdown_stmt.limit(limit_docs)
        breakdown_rows = (await session.execute(breakdown_stmt)).scalars().all()
        for doc in breakdown_rows:
            source_images = []
            for field_name in (
                "image_urls",
                "image_urls_problem_intro",
                "image_urls_causes",
                "image_urls_evaluation",
                "image_urls_inspection",
                "image_urls_solutions",
                "image_urls_key_points",
            ):
                source_images.extend(normalize_image_list(getattr(doc, field_name, None)))
            parts = [
                f"标题：{doc.title or ''}",
                f"问题描述：{doc.problem_intro or ''}",
                f"原因分析：{doc.causes or ''}",
                f"评估建议：{doc.evaluation or ''}",
                f"检查步骤：{doc.inspection or ''}",
                f"解决方案：{doc.solutions or ''}",
                f"关键要点：{doc.key_points or ''}",
            ]
            docs.append(
                {
                    "source_prefix": "db-breakdown",
                    "source_title": doc.title or f"breakdown:{doc.id}",
                    "source_doc_id": int(doc.id),
                    "source_library_type": "breakdown",
                    "source_file": doc.origin_file_dir or doc.origin_file_name or "",
                    "source_images": source_images,
                    "text": "\n".join(parts),
                }
            )

        remaining = None if not limit_docs else max(0, limit_docs - len(docs))
        if remaining != 0:
            knowledge_stmt = select(DocumentKnowledge).where(DocumentKnowledge.is_deleted == 0)
            if only_vectorized:
                knowledge_stmt = knowledge_stmt.where(DocumentKnowledge.is_vectorized == 1)
            knowledge_stmt = knowledge_stmt.order_by(DocumentKnowledge.id.asc())
            if remaining:
                knowledge_stmt = knowledge_stmt.limit(remaining)
            knowledge_rows = (await session.execute(knowledge_stmt)).scalars().all()
            for doc in knowledge_rows:
                section_rows = (
                    await session.execute(
                        select(KnowledgeDocumentSection)
                        .where(KnowledgeDocumentSection.document_id == doc.id)
                        .order_by(
                            KnowledgeDocumentSection.section_index.asc(),
                            KnowledgeDocumentSection.id.asc(),
                        )
                    )
                ).scalars().all()
                section_text = "\n\n".join(
                    f"{section.section_title or '未命名章节'}：{section.plain_text or ''}"
                    for section in section_rows
                    if str(section.plain_text or "").strip()
                )
                source_images = normalize_image_list(doc.image_urls)
                for section in section_rows:
                    source_images.extend(normalize_image_list(section.image_urls))
                docs.append(
                    {
                        "source_prefix": "db-knowledge",
                        "source_title": doc.title or f"knowledge:{doc.id}",
                        "source_doc_id": int(doc.id),
                        "source_library_type": "knowledge",
                        "source_file": doc.origin_file_dir or doc.origin_file_name or "",
                        "source_images": source_images,
                        "text": f"标题：{doc.title or ''}\n{section_text}",
                    }
                )
    return docs


def collect_upload_documents(
    upload_root: Path,
    limit_docs: Optional[int],
    max_chars: int,
    extensions: set,
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    docs = []
    skipped = []
    for path in iter_files(upload_root, extensions):
        if limit_docs and len(docs) >= limit_docs:
            break
        text, error = extract_text_from_file(path, max_chars=max_chars)
        if error:
            skipped.append({"source_file": rel_path(path), "error": error})
            continue
        docs.append(
            {
                "source_prefix": "upload",
                "source_title": path.stem,
                "source_file": rel_path(path),
                "text": text,
            }
        )
    return docs, skipped


def generate_questions_for_document(
    doc: Dict[str, Any],
    *,
    questions_per_doc: int,
    base_url: str,
    api_key: str,
    model: str,
    max_tokens: int,
    temperature: float,
) -> List[Dict[str, Any]]:
    messages = build_dataset_prompt(
        doc["source_title"],
        doc.get("source_file") or f"{doc.get('source_library_type')}:{doc.get('source_doc_id')}",
        doc["text"],
        questions_per_doc,
    )
    response = chat_completion(
        messages,
        base_url=base_url,
        api_key=api_key,
        model=model,
        max_tokens=max_tokens,
        temperature=temperature,
    )
    try:
        raw_items = extract_json_array(response)
    except json.JSONDecodeError:
        repair_response = chat_completion(
            build_json_repair_prompt(response),
            base_url=base_url,
            api_key=api_key,
            model=model,
            max_tokens=max_tokens,
            temperature=0.0,
        )
        try:
            raw_items = extract_json_array(repair_response)
        except Exception as exc:
            raise ValueError(
                f"JSON parse failed after repair: {exc}; raw_response={response[:500]}"
            ) from exc
    normalized = []
    for item in raw_items:
        qa = normalize_qa_item(
            item,
            source_prefix=doc["source_prefix"],
            source_title=doc["source_title"],
            source_file=doc.get("source_file", ""),
            source_doc_id=doc.get("source_doc_id"),
            source_library_type=doc.get("source_library_type"),
            source_images=doc.get("source_images"),
        )
        if qa:
            normalized.append(qa)
    return normalized


def with_image_question_prefix(question: str) -> str:
    question = str(question or "").strip()
    if question.startswith(("结合图片", "根据图片", "请结合图片", "请根据图片")):
        return question
    return f"结合图片，{question}"


def make_image_eval_item(
    item: Dict[str, Any],
    doc: Dict[str, Any],
    query_images: List[str],
) -> Dict[str, Any]:
    qa = dict(item)
    question = with_image_question_prefix(qa.get("question", ""))
    source_key = stable_case_id(
        "image-source",
        doc.get("source_title"),
        doc.get("source_file"),
        doc.get("source_doc_id"),
        doc.get("source_library_type"),
        ",".join(query_images),
    )
    qa.update(
        {
            "id": stable_case_id(
                "image",
                doc.get("source_title"),
                doc.get("source_file"),
                doc.get("source_doc_id"),
                doc.get("source_library_type"),
                ",".join(query_images),
                question,
            ),
            "source_key": source_key,
            "question": question,
            "query_images": list(query_images),
            "source_images": existing_image_list(doc.get("source_images")),
            "image_eval_mode": "source_image_as_query",
        }
    )
    return qa


def format_references(references: List[Dict[str, Any]]) -> str:
    parts = []
    for ref in references or []:
        title = str(ref.get("title") or "").strip()
        score = ref.get("score")
        score_text = f"{float(score):.3f}" if isinstance(score, (int, float)) else ""
        parts.append(
            f"{ref.get('library_type', 'breakdown')}:{ref.get('doc_id')} {title} {score_text}".strip()
        )
    return "; ".join(parts)


def format_contexts_for_judge(references: List[Dict[str, Any]], max_chars: int = 6000) -> str:
    parts = []
    total = 0
    for ref_index, ref in enumerate(references or [], start=1):
        header = (
            f"[{ref_index}] {ref.get('library_type', 'breakdown')}:{ref.get('doc_id')} "
            f"{ref.get('title') or ''} score={ref.get('score')}"
        ).strip()
        chunks = ref.get("chunks") or []
        if chunks:
            chunk_texts = []
            for chunk_index, chunk in enumerate(chunks[:3], start=1):
                text = str(chunk.get("content") or "").strip()
                if text:
                    chunk_texts.append(f"片段{chunk_index}：{text}")
            body = "\n".join(chunk_texts)
        else:
            body = ""
        block = f"{header}\n{body}".strip()
        if not block:
            continue
        if total + len(block) > max_chars:
            remaining = max_chars - total
            if remaining <= 0:
                break
            block = block[:remaining]
        parts.append(block)
        total += len(block)
        if total >= max_chars:
            break
    return "\n\n".join(parts)


def compact_references(references: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    compact = []
    for ref in references or []:
        chunks = []
        for chunk in (ref.get("chunks") or [])[:3]:
            chunks.append(
                {
                    "content": trim_text(str(chunk.get("content") or ""), 900),
                    "score": chunk.get("score"),
                    "metadata": chunk.get("metadata") if isinstance(chunk.get("metadata"), dict) else {},
                }
            )
        compact.append(
            {
                "doc_id": ref.get("doc_id"),
                "library_type": ref.get("library_type"),
                "title": ref.get("title"),
                "score": ref.get("score"),
                "chunks": chunks,
            }
        )
    return compact


def references_from_backend_payload(payload: Any) -> List[Dict[str, Any]]:
    if isinstance(payload, dict) and "data" in payload:
        payload = payload.get("data")
    if not isinstance(payload, list):
        return []
    return compact_references([item for item in payload if isinstance(item, dict)])


def normalize_relevant_docs(item: Dict[str, Any]) -> List[Tuple[str, int]]:
    candidates = item.get("relevant_docs") or item.get("expected_docs") or item.get("ground_truth_docs")
    relevant: List[Tuple[str, int]] = []

    if isinstance(candidates, list):
        for candidate in candidates:
            if not isinstance(candidate, dict):
                continue
            doc_id = candidate.get("doc_id") or candidate.get("source_doc_id")
            library_type = candidate.get("library_type") or candidate.get("source_library_type")
            if doc_id is None or not library_type:
                continue
            try:
                relevant.append((str(library_type), int(doc_id)))
            except Exception:
                continue

    if not relevant:
        doc_id = item.get("source_doc_id")
        library_type = item.get("source_library_type")
        if doc_id is not None and library_type:
            try:
                relevant.append((str(library_type), int(doc_id)))
            except Exception:
                pass

    deduped = []
    seen = set()
    for value in relevant:
        if value in seen:
            continue
        deduped.append(value)
        seen.add(value)
    return deduped


def reference_key(ref: Dict[str, Any]) -> Optional[Tuple[str, int]]:
    doc_id = ref.get("doc_id")
    library_type = ref.get("library_type")
    if doc_id is None or not library_type:
        return None
    try:
        return str(library_type), int(doc_id)
    except Exception:
        return None


def normalize_backend_url(url: str) -> str:
    url = (url or "").strip().rstrip("/")
    if url.endswith("/api/v1"):
        return url
    return url + "/api/v1"


def expected_doc_rank(item: Dict[str, Any], references: List[Dict[str, Any]]) -> Optional[int]:
    relevant_docs = set(normalize_relevant_docs(item))
    if not relevant_docs:
        return None
    for index, ref in enumerate(references or [], start=1):
        if reference_key(ref) in relevant_docs:
            return index
    return 0


def expected_doc_hit(item: Dict[str, Any], references: List[Dict[str, Any]]) -> Optional[bool]:
    rank = expected_doc_rank(item, references)
    if rank is None:
        return None
    return rank > 0


def retrieval_metrics_for_row(row: Dict[str, Any]) -> Optional[Dict[str, float]]:
    relevant_docs = set(normalize_relevant_docs(row))
    if not relevant_docs:
        return None

    references = row.get("references") or []
    k = len(references)
    if k <= 0:
        return {
            "precision": 0.0,
            "recall": 0.0,
            "f1": 0.0,
            "reciprocal_rank": 0.0,
            "average_precision": 0.0,
        }

    hit_count = 0
    reciprocal_rank = 0.0
    precision_sum = 0.0
    seen_relevant = set()
    for index, ref in enumerate(references, start=1):
        key = reference_key(ref)
        if key not in relevant_docs or key in seen_relevant:
            continue
        seen_relevant.add(key)
        hit_count += 1
        precision_at_rank = hit_count / index
        precision_sum += precision_at_rank
        if reciprocal_rank == 0.0:
            reciprocal_rank = 1 / index

    precision = hit_count / k
    recall = hit_count / len(relevant_docs)
    f1 = (2 * precision * recall / (precision + recall)) if (precision + recall) else 0.0
    average_precision = precision_sum / len(relevant_docs)
    return {
        "precision": precision,
        "recall": recall,
        "f1": f1,
        "reciprocal_rank": reciprocal_rank,
        "average_precision": average_precision,
    }


async def retrieve_current_rag(question: str, query_images: str = "") -> Dict[str, Any]:
    from database import AsyncSessionLocal
    from routers.message import get_reference_documents

    started = time.perf_counter()
    async with AsyncSessionLocal() as session:
        references = await get_reference_documents(session, question, query_images or None)

    return {
        "references": compact_references(references),
        "query_images": query_images,
        "latency_seconds": round(time.perf_counter() - started, 3),
    }


def login_backend(
    *,
    backend_url: str,
    username: str,
    password: str,
    role: str = "",
) -> str:
    import requests

    payload = {"username": username, "password": password}
    if role:
        payload["role"] = role
    response = requests.post(
        normalize_backend_url(backend_url) + "/auth/login",
        json=payload,
        timeout=60,
    )
    response.raise_for_status()
    body = response.json()
    token = ((body.get("data") or {}).get("access_token") or "").strip()
    if not token:
        raise RuntimeError(f"login succeeded but no access_token returned: {body}")
    return token


def create_backend_conversation(*, backend_url: str, token: str) -> int:
    import requests

    response = requests.post(
        normalize_backend_url(backend_url) + "/conversation/create",
        headers={"Authorization": f"Bearer {token}"},
        timeout=60,
    )
    response.raise_for_status()
    body = response.json()
    session_id = (body.get("data") or {}).get("id")
    if session_id is None:
        raise RuntimeError(f"create conversation returned no id: {body}")
    return int(session_id)


def ask_backend_stream(
    *,
    backend_url: str,
    token: str,
    session_id: int,
    question: str,
    query_images: str = "",
    stop_after_references: bool = False,
) -> Dict[str, Any]:
    import requests

    started = time.perf_counter()
    response = requests.post(
        normalize_backend_url(backend_url) + "/message/ask",
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
        },
        json={
            "session_id": session_id,
            "content_text": question,
            "user_uploaded_images": query_images or None,
            "stream": True,
        },
        stream=True,
        timeout=300,
    )
    response.raise_for_status()

    answer = ""
    references: List[Dict[str, Any]] = []
    for line in response.iter_lines(decode_unicode=True):
        if not line:
            continue
        line = line.strip()
        if line.startswith("data:"):
            line = line[5:].strip()
        if not line:
            continue
        try:
            payload = json.loads(line)
        except json.JSONDecodeError:
            continue
        if payload.get("code") == 1 and payload.get("data") == "true":
            break
        if payload.get("reference_docs"):
            references = compact_references(payload.get("reference_docs") or [])
            if stop_after_references:
                break
        if isinstance(payload.get("answer"), str):
            answer = payload.get("answer") or answer

    return {
        "answer": answer.strip(),
        "references": references,
        "query_images": query_images,
        "latency_seconds": round(time.perf_counter() - started, 3),
    }


async def ask_current_rag(
    question: str,
    *,
    query_images: str = "",
    base_url: str,
    api_key: str,
    model: str,
    max_tokens: int,
) -> Dict[str, Any]:
    from database import AsyncSessionLocal
    from routers.message import (
        generate_messages,
        get_ai_reference_prompt_refs,
        get_reference_documents,
    )

    started = time.perf_counter()
    async with AsyncSessionLocal() as session:
        references = await get_reference_documents(session, question, query_images or None)
        prompt_refs = get_ai_reference_prompt_refs(references)
        fake_message = SimpleNamespace(
            content_text=question,
            user_uploaded_images=query_images or None,
            message_order=1,
            token_count=0,
        )
        messages = await generate_messages(session, 0, fake_message, prompt_refs)

    answer = await asyncio.to_thread(
        chat_completion,
        messages,
        base_url=base_url,
        api_key=api_key,
        model=model,
        max_tokens=max_tokens,
        temperature=0.0,
    )
    return {
        "answer": answer.strip(),
        "references": compact_references(references),
        "query_images": query_images,
        "latency_seconds": round(time.perf_counter() - started, 3),
    }


def judge_answer(
    question: str,
    ground_truth: str,
    answer: str,
    *,
    base_url: str,
    api_key: str,
    model: str,
    max_tokens: int,
    threshold: float,
) -> Dict[str, Any]:
    if not answer.strip():
        return {"score": 0.0, "passed": False, "reason": "模型答案为空"}
    messages = build_judge_prompt(question, ground_truth, answer)
    response = chat_completion(
        messages,
        base_url=base_url,
        api_key=api_key,
        model=model,
        max_tokens=max_tokens,
        temperature=0.0,
    )
    try:
        payload = extract_json_object(response)
        score = float(payload.get("score", 0.0))
        score = max(0.0, min(1.0, score))
        passed = bool(payload.get("passed", score >= threshold))
        return {
            "score": score,
            "passed": passed and score >= threshold,
            "reason": str(payload.get("reason") or "").strip(),
        }
    except Exception:
        return {
            "score": 0.0,
            "passed": False,
            "reason": f"裁判输出解析失败：{response[:200]}",
        }


def judge_rag_metrics(
    question: str,
    ground_truth: str,
    answer: str,
    references: List[Dict[str, Any]],
    *,
    base_url: str,
    api_key: str,
    model: str,
    max_tokens: int,
) -> Dict[str, Any]:
    if not answer.strip():
        result = {name: 0.0 for name in RAG_METRIC_NAMES}
        result["reason"] = "模型答案为空"
        return result

    contexts = format_contexts_for_judge(references)
    messages = build_rag_metrics_prompt(question, ground_truth, answer, contexts)
    response = chat_completion(
        messages,
        base_url=base_url,
        api_key=api_key,
        model=model,
        max_tokens=max_tokens,
        temperature=0.0,
    )
    try:
        payload = extract_json_object(response)
        result = {}
        for name in RAG_METRIC_NAMES:
            score = float(payload.get(name, 0.0))
            result[name] = round(max(0.0, min(1.0, score)), 4)
        result["reason"] = str(payload.get("reason") or "").strip()
        return result
    except Exception:
        result = {name: 0.0 for name in RAG_METRIC_NAMES}
        result["reason"] = f"RAG指标裁判输出解析失败：{response[:200]}"
        return result


def compute_summary(rows: List[Dict[str, Any]], judge_enabled: bool) -> Dict[str, Any]:
    total = len(rows)
    answered = sum(1 for row in rows if str(row.get("answer") or "").strip())
    summary = {
        "total_cases": total,
        "answered_cases": answered,
        "answer_rate": round(answered / total, 4) if total else 0,
    }
    retrieval_rows = [
        row for row in rows
        if normalize_relevant_docs(row)
    ]
    if retrieval_rows:
        ranks = []
        metric_rows = []
        for row in retrieval_rows:
            rank = row.get("retrieval_rank")
            if rank is None:
                rank = expected_doc_rank(row, row.get("references") or [])
            try:
                ranks.append(int(rank))
            except Exception:
                ranks.append(0)
            metrics = retrieval_metrics_for_row(row)
            if metrics:
                metric_rows.append(metrics)
        hit_ranks = [rank for rank in ranks if rank > 0]
        summary["retrieval_cases"] = len(retrieval_rows)
        summary["retrieval_top1"] = round(sum(1 for rank in ranks if rank == 1) / len(ranks), 4)
        summary["retrieval_topk"] = round(len(hit_ranks) / len(ranks), 4)
        summary["retrieval_mrr"] = round(sum((1 / rank) if rank > 0 else 0 for rank in ranks) / len(ranks), 4)
        summary["retrieval_misses"] = sum(1 for rank in ranks if rank <= 0)
        summary["average_retrieved_docs"] = round(
            sum(len(row.get("references") or []) for row in retrieval_rows) / len(retrieval_rows),
            4,
        )
        if metric_rows:
            summary["precision_at_k"] = round(
                sum(item["precision"] for item in metric_rows) / len(metric_rows),
                4,
            )
            summary["recall_at_k"] = round(
                sum(item["recall"] for item in metric_rows) / len(metric_rows),
                4,
            )
            summary["f1_at_k"] = round(
                sum(item["f1"] for item in metric_rows) / len(metric_rows),
                4,
            )
            summary["mrr"] = round(
                sum(item["reciprocal_rank"] for item in metric_rows) / len(metric_rows),
                4,
            )
            summary["map"] = round(
                sum(item["average_precision"] for item in metric_rows) / len(metric_rows),
                4,
            )
    if judge_enabled:
        scores = [float(row.get("score", 0.0)) for row in rows]
        passed = [bool(row.get("passed", False)) for row in rows]
        summary["accuracy"] = round(sum(1 for value in passed if value) / total, 4) if total else 0
        summary["average_score"] = round(sum(scores) / total, 4) if total else 0
        summary["passed_cases"] = sum(1 for value in passed if value)
        for metric_name in RAG_METRIC_NAMES:
            metric_scores = [
                float(row.get(metric_name))
                for row in rows
                if row.get(metric_name) is not None
            ]
            summary[metric_name] = (
                round(sum(metric_scores) / len(metric_scores), 4)
                if metric_scores
                else 0
            )
    return summary


def export_excel(rows: List[Dict[str, Any]], summary: Dict[str, Any], path: Path) -> None:
    import pandas as pd

    ensure_parent(path)
    report_rows = []
    for row in rows:
        retrieval_metrics = retrieval_metrics_for_row(row) or {}
        report_rows.append(
            {
                "id": row.get("id"),
                "question": row.get("question"),
                "query_images": row.get("query_images"),
                "source_images": ", ".join(normalize_image_list(row.get("source_images"))),
                "ground_truth": row.get("ground_truth"),
                "answer": row.get("answer"),
                "context_precision": row.get("context_precision"),
                "context_recall": row.get("context_recall"),
                "faithfulness": row.get("faithfulness"),
                "answer_accuracy": row.get("answer_accuracy"),
                "score": row.get("score"),
                "passed": row.get("passed"),
                "judge_reason": row.get("judge_reason"),
                "rag_metric_reason": row.get("rag_metric_reason"),
                "retrieval_rank": row.get("retrieval_rank"),
                "retrieved_expected_doc": row.get("retrieved_expected_doc"),
                "precision_at_k": (
                    round(retrieval_metrics["precision"], 4)
                    if "precision" in retrieval_metrics
                    else None
                ),
                "recall_at_k": (
                    round(retrieval_metrics["recall"], 4)
                    if "recall" in retrieval_metrics
                    else None
                ),
                "f1_at_k": (
                    round(retrieval_metrics["f1"], 4)
                    if "f1" in retrieval_metrics
                    else None
                ),
                "average_precision": (
                    round(retrieval_metrics["average_precision"], 4)
                    if "average_precision" in retrieval_metrics
                    else None
                ),
                "references": format_references(row.get("references") or []),
                "latency_seconds": row.get("latency_seconds"),
                "source_title": row.get("source_title"),
                "source_file": row.get("source_file"),
                "source_library_type": row.get("source_library_type"),
                "source_doc_id": row.get("source_doc_id"),
                "error": row.get("error"),
            }
        )
    summary_rows = [{"metric": key, "value": value} for key, value in summary.items()]
    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        pd.DataFrame(report_rows).to_excel(writer, index=False, sheet_name="cases")
        pd.DataFrame(summary_rows).to_excel(writer, index=False, sheet_name="summary")


async def command_inventory(args: argparse.Namespace) -> None:
    upload_root = Path(args.upload_root)
    if not upload_root.is_absolute():
        upload_root = PROJECT_ROOT / upload_root

    payload = {"upload": file_inventory(upload_root)}
    try:
        payload["database"] = await db_inventory()
    except Exception as exc:
        payload["database_error"] = f"{type(exc).__name__}: {exc}"

    if args.out:
        write_json(Path(args.out), payload)
    print(json.dumps(payload, ensure_ascii=False, indent=2))


async def command_make_dataset(args: argparse.Namespace) -> None:
    answer_base_url, api_key, model = resolve_ai_config(args.base_url, args.api_key, args.model)
    out_path = Path(args.out)
    existing = load_json_list(out_path)
    seen_ids = {item.get("id") for item in existing}
    seen_source_keys = {item.get("source_key") for item in existing if item.get("source_key")}
    dataset = list(existing)

    if args.source == "db":
        docs = await collect_db_documents(args.limit_docs, only_vectorized=not args.include_unvectorized)
        skipped = []
    else:
        upload_root = Path(args.upload_root)
        if not upload_root.is_absolute():
            upload_root = PROJECT_ROOT / upload_root
        docs, skipped = collect_upload_documents(
            upload_root,
            limit_docs=args.limit_docs,
            max_chars=args.max_chars,
            extensions=normalize_extensions(args.extensions),
        )

    skipped_path = out_path.with_suffix(".skipped.json")
    if skipped:
        write_json(skipped_path, skipped)

    print(f"documents: {len(docs)}, existing questions: {len(existing)}")
    for index, doc in enumerate(docs, start=1):
        doc_key = stable_case_id(
            doc["source_prefix"],
            doc.get("source_title"),
            doc.get("source_file"),
            doc.get("source_doc_id"),
            doc.get("source_library_type"),
        )
        if doc_key in seen_source_keys:
            continue
        print(f"[{index}/{len(docs)}] generating questions: {doc['source_title']}")
        try:
            qa_items = generate_questions_for_document(
                doc,
                questions_per_doc=args.questions_per_doc,
                base_url=answer_base_url,
                api_key=api_key,
                model=model,
                max_tokens=args.max_tokens,
                temperature=args.temperature,
            )
            new_items = [item for item in qa_items if item.get("id") not in seen_ids]
            dataset.extend(new_items)
            seen_ids.update(item["id"] for item in new_items)
            seen_source_keys.update(item.get("source_key") for item in new_items if item.get("source_key"))
            write_json(out_path, dataset)
            print(f"  added: {len(new_items)}, total: {len(dataset)}")
        except Exception as exc:
            skipped.append(
                {
                    "source_title": doc.get("source_title"),
                    "source_file": doc.get("source_file"),
                    "error": f"{type(exc).__name__}: {exc}",
                }
            )
            write_json(skipped_path, skipped)
            print(f"  failed: {type(exc).__name__}: {exc}")

    print(f"dataset saved: {out_path}")
    if skipped:
        print(f"skipped/failed saved: {skipped_path}")


async def command_make_image_dataset(args: argparse.Namespace) -> None:
    answer_base_url, api_key, model = resolve_ai_config(args.base_url, args.api_key, args.model)
    out_path = Path(args.out)
    image_dir = Path(args.image_dir)
    if not image_dir.is_absolute():
        image_dir = PROJECT_ROOT / image_dir
    existing = load_json_list(out_path)
    seen_ids = {item.get("id") for item in existing}
    seen_source_keys = {item.get("source_key") for item in existing if item.get("source_key")}
    dataset = list(existing)
    skipped: List[Dict[str, Any]] = []

    docs = await collect_db_documents(args.limit_docs, only_vectorized=not args.include_unvectorized)
    image_docs = []
    for doc in docs:
        images = existing_image_list(doc.get("source_images"))
        if not images:
            skipped.append(
                {
                    "source_title": doc.get("source_title"),
                    "source_file": doc.get("source_file"),
                    "reason": "no existing source_images",
                }
            )
            continue
        doc = dict(doc)
        doc["source_images"] = images
        image_docs.append(doc)

    skipped_path = out_path.with_suffix(".skipped.json")
    if skipped:
        write_json(skipped_path, skipped)

    print(
        f"image documents: {len(image_docs)}, skipped without images: {len(skipped)}, "
        f"existing questions: {len(existing)}"
    )
    for index, doc in enumerate(image_docs, start=1):
        source_images = existing_image_list(doc.get("source_images"))
        case_seed = stable_case_id(
            "image-copy",
            doc.get("source_title"),
            doc.get("source_file"),
            doc.get("source_doc_id"),
            doc.get("source_library_type"),
        )
        query_images = copy_eval_images(
            source_images,
            image_dir,
            case_seed,
            args.max_images_per_case,
        )
        doc_key = stable_case_id(
            "image-source",
            doc.get("source_title"),
            doc.get("source_file"),
            doc.get("source_doc_id"),
            doc.get("source_library_type"),
            ",".join(query_images),
        )
        if doc_key in seen_source_keys:
            continue
        print(f"[{index}/{len(image_docs)}] generating image questions: {doc['source_title']}")
        try:
            qa_items = generate_questions_for_document(
                doc,
                questions_per_doc=args.questions_per_doc,
                base_url=answer_base_url,
                api_key=api_key,
                model=model,
                max_tokens=args.max_tokens,
                temperature=args.temperature,
            )
            image_items = [make_image_eval_item(item, doc, query_images) for item in qa_items]
            new_items = [item for item in image_items if item.get("id") not in seen_ids]
            dataset.extend(new_items)
            seen_ids.update(item["id"] for item in new_items)
            seen_source_keys.update(item.get("source_key") for item in new_items if item.get("source_key"))
            write_json(out_path, dataset)
            print(f"  images: {len(query_images)}, added: {len(new_items)}, total: {len(dataset)}")
        except Exception as exc:
            skipped.append(
                {
                    "source_title": doc.get("source_title"),
                    "source_file": doc.get("source_file"),
                    "source_images": query_images,
                    "error": f"{type(exc).__name__}: {exc}",
                }
            )
            write_json(skipped_path, skipped)
            print(f"  failed: {type(exc).__name__}: {exc}")

    print(f"image dataset saved: {out_path}")
    print(f"image folder: {image_dir}")
    if skipped:
        print(f"skipped/failed saved: {skipped_path}")


async def command_retrieve(args: argparse.Namespace) -> None:
    dataset_path = Path(args.dataset)
    out_dir = Path(args.out_dir) if args.out_dir else PROJECT_ROOT / "evaluate" / f"retrieval_eval_{now_stamp()}"
    out_dir.mkdir(parents=True, exist_ok=True)

    result_path = out_dir / "retrieval_results.json"
    summary_path = out_dir / "retrieval_summary.json"
    debug_log_path = out_dir / "retrieval_debug.log"
    excel_path = out_dir / "retrieval_accuracy_report.xlsx"

    dataset = load_json_list(dataset_path)
    if args.limit_cases:
        dataset = dataset[: args.limit_cases]

    rows = load_json_list(result_path)
    completed_ids = {row.get("id") for row in rows}

    print(f"retrieval cases: {len(dataset)}, completed: {len(completed_ids)}, out: {out_dir}")
    for index, item in enumerate(dataset, start=1):
        case_id = item.get("id") or stable_case_id("case", index, item.get("question"))
        if case_id in completed_ids:
            continue

        question = str(item.get("question") or "").strip()
        if not question:
            continue
        query_images = select_query_images(
            item,
            use_source_images=args.use_source_images,
            max_query_images=args.max_query_images,
        )

        print(f"[{index}/{len(dataset)}] {question[:80]}")
        debug_buffer = io.StringIO()
        try:
            if args.show_rag_debug:
                payload = await retrieve_current_rag(question, query_images)
            else:
                with contextlib.redirect_stdout(debug_buffer):
                    payload = await retrieve_current_rag(question, query_images)
                with debug_log_path.open("a", encoding="utf-8") as f:
                    f.write(f"\n\n===== {case_id} {question} =====\n")
                    f.write(debug_buffer.getvalue())

            row = dict(item)
            row["id"] = case_id
            row["answer"] = ""
            row.update(payload)
            row["retrieval_rank"] = expected_doc_rank(row, row.get("references") or [])
            row["retrieved_expected_doc"] = expected_doc_hit(row, row.get("references") or [])
            rows.append(row)
            completed_ids.add(case_id)
            write_json(result_path, rows)
            summary = compute_summary(rows, judge_enabled=False)
            write_json(summary_path, summary)
            print(
                f"  rank={row.get('retrieval_rank')} "
                f"hit={row.get('retrieved_expected_doc')} refs={len(row.get('references') or [])}"
            )
        except Exception as exc:
            row = dict(item)
            row["id"] = case_id
            row["answer"] = ""
            row["references"] = []
            row["retrieval_rank"] = 0
            row["retrieved_expected_doc"] = False
            row["error"] = f"{type(exc).__name__}: {exc}"
            rows.append(row)
            completed_ids.add(case_id)
            write_json(result_path, rows)
            print(f"  failed: {row['error']}")

    summary = compute_summary(rows, judge_enabled=False)
    write_json(summary_path, summary)
    try:
        export_excel(rows, summary, excel_path)
    except Exception as exc:
        print(f"excel export failed: {type(exc).__name__}: {exc}")

    print(json.dumps(summary, ensure_ascii=False, indent=2))
    print(f"results: {result_path}")
    print(f"summary: {summary_path}")
    print(f"excel: {excel_path}")
    if not args.show_rag_debug and debug_log_path.exists():
        print(f"rag debug log: {debug_log_path}")


async def command_run(args: argparse.Namespace) -> None:
    dataset_path = Path(args.dataset)
    out_dir = Path(args.out_dir) if args.out_dir else PROJECT_ROOT / "evaluate" / f"company_demo_{now_stamp()}"
    out_dir.mkdir(parents=True, exist_ok=True)

    result_path = out_dir / "results.json"
    summary_path = out_dir / "summary.json"
    debug_log_path = out_dir / "rag_debug.log"
    excel_path = out_dir / "qa_accuracy_report.xlsx"

    answer_base_url, api_key, answer_model = resolve_ai_config(
        args.answer_base_url,
        args.api_key,
        args.answer_model,
    )
    judge_base_url = args.judge_base_url or os.getenv("JUDGE_AI_BASE_URL") or answer_base_url
    judge_model = args.judge_model or os.getenv("JUDGE_MODEL") or answer_model

    dataset = load_json_list(dataset_path)
    if args.limit_cases:
        dataset = dataset[: args.limit_cases]

    rows = load_json_list(result_path)
    completed_ids = {row.get("id") for row in rows}

    judge_enabled = not args.no_judge
    print(f"cases: {len(dataset)}, completed: {len(completed_ids)}, out: {out_dir}")

    for index, item in enumerate(dataset, start=1):
        case_id = item.get("id") or stable_case_id("case", index, item.get("question"))
        if case_id in completed_ids:
            continue

        question = str(item.get("question") or "").strip()
        ground_truth = str(item.get("ground_truth") or "").strip()
        if not question or not ground_truth:
            continue
        query_images = select_query_images(
            item,
            use_source_images=args.use_source_images,
            max_query_images=args.max_query_images,
        )

        print(f"[{index}/{len(dataset)}] {question[:80]}")
        debug_buffer = io.StringIO()
        try:
            if args.show_rag_debug:
                rag_payload = await ask_current_rag(
                    question,
                    query_images=query_images,
                    base_url=answer_base_url,
                    api_key=api_key,
                    model=answer_model,
                    max_tokens=args.max_answer_tokens,
                )
            else:
                with contextlib.redirect_stdout(debug_buffer):
                    rag_payload = await ask_current_rag(
                        question,
                        query_images=query_images,
                        base_url=answer_base_url,
                        api_key=api_key,
                        model=answer_model,
                        max_tokens=args.max_answer_tokens,
                    )
                with debug_log_path.open("a", encoding="utf-8") as f:
                    f.write(f"\n\n===== {case_id} {question} =====\n")
                    f.write(debug_buffer.getvalue())

            row = dict(item)
            row["id"] = case_id
            row.update(rag_payload)
            row["retrieval_rank"] = expected_doc_rank(row, row.get("references") or [])
            row["retrieved_expected_doc"] = expected_doc_hit(row, row.get("references") or [])

            if judge_enabled:
                rag_metrics = await asyncio.to_thread(
                    judge_rag_metrics,
                    question,
                    ground_truth,
                    row.get("answer") or "",
                    row.get("references") or [],
                    base_url=judge_base_url,
                    api_key=api_key,
                    model=judge_model,
                    max_tokens=args.max_judge_tokens,
                )
                for metric_name in RAG_METRIC_NAMES:
                    row[metric_name] = rag_metrics.get(metric_name)
                row["rag_metric_reason"] = rag_metrics.get("reason")
                row["score"] = row.get("answer_accuracy", 0.0)
                row["passed"] = float(row["score"]) >= args.judge_threshold
                row["judge_reason"] = row["rag_metric_reason"]

            rows.append(row)
            completed_ids.add(case_id)
            write_json(result_path, rows)
            summary = compute_summary(rows, judge_enabled=judge_enabled)
            write_json(summary_path, summary)
            print(
                f"  cp={row.get('context_precision', 'n/a')} "
                f"cr={row.get('context_recall', 'n/a')} "
                f"faith={row.get('faithfulness', 'n/a')} "
                f"acc={row.get('answer_accuracy', row.get('score', 'n/a'))} "
                f"refs={len(row.get('references') or [])}"
            )
        except Exception as exc:
            row = dict(item)
            row["id"] = case_id
            row["answer"] = ""
            row["references"] = []
            for metric_name in RAG_METRIC_NAMES:
                row[metric_name] = 0.0 if judge_enabled else None
            row["score"] = 0.0 if judge_enabled else None
            row["passed"] = False if judge_enabled else None
            row["error"] = f"{type(exc).__name__}: {exc}"
            rows.append(row)
            completed_ids.add(case_id)
            write_json(result_path, rows)
            print(f"  failed: {row['error']}")

    summary = compute_summary(rows, judge_enabled=judge_enabled)
    write_json(summary_path, summary)
    try:
        export_excel(rows, summary, excel_path)
    except Exception as exc:
        print(f"excel export failed: {type(exc).__name__}: {exc}")

    print(json.dumps(summary, ensure_ascii=False, indent=2))
    print(f"results: {result_path}")
    print(f"summary: {summary_path}")
    print(f"excel: {excel_path}")
    if not args.show_rag_debug and debug_log_path.exists():
        print(f"rag debug log: {debug_log_path}")


async def command_retrieve_api(args: argparse.Namespace) -> None:
    dataset_path = Path(args.dataset)
    out_dir = Path(args.out_dir) if args.out_dir else PROJECT_ROOT / "evaluate" / f"retrieval_api_eval_{now_stamp()}"
    out_dir.mkdir(parents=True, exist_ok=True)

    result_path = out_dir / "retrieval_api_results.json"
    summary_path = out_dir / "retrieval_api_summary.json"
    excel_path = out_dir / "retrieval_api_report.xlsx"

    token = args.token.strip()
    if not token:
        if not args.username or not args.password:
            raise ValueError("retrieve-api 需要 --token，或者 --username 和 --password")
        token = await asyncio.to_thread(
            login_backend,
            backend_url=args.backend_url,
            username=args.username,
            password=args.password,
            role=args.role,
        )

    dataset = load_json_list(dataset_path)
    if args.limit_cases:
        dataset = dataset[: args.limit_cases]

    rows = load_json_list(result_path)
    completed_ids = {row.get("id") for row in rows}

    print(f"retrieval api cases: {len(dataset)}, completed: {len(completed_ids)}, out: {out_dir}")
    for index, item in enumerate(dataset, start=1):
        case_id = item.get("id") or stable_case_id("case", index, item.get("question"))
        if case_id in completed_ids:
            continue

        question = str(item.get("question") or "").strip()
        if not question:
            continue
        query_images = select_query_images(
            item,
            use_source_images=args.use_source_images,
            max_query_images=args.max_query_images,
        )

        print(f"[{index}/{len(dataset)}] {question[:80]}")
        try:
            session_id = await asyncio.to_thread(
                create_backend_conversation,
                backend_url=args.backend_url,
                token=token,
            )
            api_payload = await asyncio.to_thread(
                ask_backend_stream,
                backend_url=args.backend_url,
                token=token,
                session_id=session_id,
                question=question,
                query_images=query_images,
                stop_after_references=True,
            )

            row = dict(item)
            row["id"] = case_id
            row["answer"] = ""
            row["references"] = api_payload.get("references") or []
            row["latency_seconds"] = api_payload.get("latency_seconds")
            row["retrieval_rank"] = expected_doc_rank(row, row.get("references") or [])
            row["retrieved_expected_doc"] = expected_doc_hit(row, row.get("references") or [])
            rows.append(row)
            completed_ids.add(case_id)
            write_json(result_path, rows)
            summary = compute_summary(rows, judge_enabled=False)
            write_json(summary_path, summary)
            print(
                f"  rank={row.get('retrieval_rank')} "
                f"hit={row.get('retrieved_expected_doc')} refs={len(row.get('references') or [])}"
            )
        except Exception as exc:
            row = dict(item)
            row["id"] = case_id
            row["answer"] = ""
            row["references"] = []
            row["retrieval_rank"] = 0
            row["retrieved_expected_doc"] = False
            row["error"] = f"{type(exc).__name__}: {exc}"
            rows.append(row)
            completed_ids.add(case_id)
            write_json(result_path, rows)
            print(f"  failed: {row['error']}")

    summary = compute_summary(rows, judge_enabled=False)
    write_json(summary_path, summary)
    try:
        export_excel(rows, summary, excel_path)
    except Exception as exc:
        print(f"excel export failed: {type(exc).__name__}: {exc}")

    print(json.dumps(summary, ensure_ascii=False, indent=2))
    print(f"results: {result_path}")
    print(f"summary: {summary_path}")
    print(f"excel: {excel_path}")


async def command_run_api(args: argparse.Namespace) -> None:
    dataset_path = Path(args.dataset)
    out_dir = Path(args.out_dir) if args.out_dir else PROJECT_ROOT / "evaluate" / f"api_eval_{now_stamp()}"
    out_dir.mkdir(parents=True, exist_ok=True)

    result_path = out_dir / "api_results.json"
    summary_path = out_dir / "api_summary.json"
    excel_path = out_dir / "api_accuracy_report.xlsx"

    answer_base_url, api_key, answer_model = resolve_ai_config(
        "",
        args.api_key,
        args.judge_model,
    )
    judge_base_url = args.judge_base_url or os.getenv("JUDGE_AI_BASE_URL") or answer_base_url
    judge_model = args.judge_model or os.getenv("JUDGE_MODEL") or answer_model

    token = args.token.strip()
    if not token:
        if not args.username or not args.password:
            raise ValueError("run-api 需要 --token，或者 --username 和 --password")
        token = await asyncio.to_thread(
            login_backend,
            backend_url=args.backend_url,
            username=args.username,
            password=args.password,
            role=args.role,
        )

    dataset = load_json_list(dataset_path)
    if args.limit_cases:
        dataset = dataset[: args.limit_cases]

    rows = load_json_list(result_path)
    completed_ids = {row.get("id") for row in rows}
    judge_enabled = not args.no_judge

    print(f"api cases: {len(dataset)}, completed: {len(completed_ids)}, out: {out_dir}")
    for index, item in enumerate(dataset, start=1):
        case_id = item.get("id") or stable_case_id("case", index, item.get("question"))
        if case_id in completed_ids:
            continue

        question = str(item.get("question") or "").strip()
        ground_truth = str(item.get("ground_truth") or "").strip()
        if not question or not ground_truth:
            continue
        query_images = select_query_images(
            item,
            use_source_images=args.use_source_images,
            max_query_images=args.max_query_images,
        )

        print(f"[{index}/{len(dataset)}] {question[:80]}")
        try:
            session_id = await asyncio.to_thread(
                create_backend_conversation,
                backend_url=args.backend_url,
                token=token,
            )
            api_payload = await asyncio.to_thread(
                ask_backend_stream,
                backend_url=args.backend_url,
                token=token,
                session_id=session_id,
                question=question,
                query_images=query_images,
            )

            row = dict(item)
            row["id"] = case_id
            row.update(api_payload)
            row["retrieval_rank"] = expected_doc_rank(row, row.get("references") or [])
            row["retrieved_expected_doc"] = expected_doc_hit(row, row.get("references") or [])

            if judge_enabled:
                rag_metrics = await asyncio.to_thread(
                    judge_rag_metrics,
                    question,
                    ground_truth,
                    row.get("answer") or "",
                    row.get("references") or [],
                    base_url=judge_base_url,
                    api_key=api_key,
                    model=judge_model,
                    max_tokens=args.max_judge_tokens,
                )
                for metric_name in RAG_METRIC_NAMES:
                    row[metric_name] = rag_metrics.get(metric_name)
                row["rag_metric_reason"] = rag_metrics.get("reason")
                row["score"] = row.get("answer_accuracy", 0.0)
                row["passed"] = float(row["score"]) >= args.judge_threshold
                row["judge_reason"] = row["rag_metric_reason"]

            rows.append(row)
            completed_ids.add(case_id)
            write_json(result_path, rows)
            summary = compute_summary(rows, judge_enabled=judge_enabled)
            write_json(summary_path, summary)
            print(
                f"  rank={row.get('retrieval_rank')} "
                f"cp={row.get('context_precision', 'n/a')} "
                f"cr={row.get('context_recall', 'n/a')} "
                f"faith={row.get('faithfulness', 'n/a')} "
                f"acc={row.get('answer_accuracy', row.get('score', 'n/a'))} "
                f"refs={len(row.get('references') or [])}"
            )
        except Exception as exc:
            row = dict(item)
            row["id"] = case_id
            row["answer"] = ""
            row["references"] = []
            row["retrieval_rank"] = 0
            row["retrieved_expected_doc"] = False
            for metric_name in RAG_METRIC_NAMES:
                row[metric_name] = 0.0 if judge_enabled else None
            row["score"] = 0.0 if judge_enabled else None
            row["passed"] = False if judge_enabled else None
            row["error"] = f"{type(exc).__name__}: {exc}"
            rows.append(row)
            completed_ids.add(case_id)
            write_json(result_path, rows)
            print(f"  failed: {row['error']}")

    summary = compute_summary(rows, judge_enabled=judge_enabled)
    write_json(summary_path, summary)
    try:
        export_excel(rows, summary, excel_path)
    except Exception as exc:
        print(f"excel export failed: {type(exc).__name__}: {exc}")

    print(json.dumps(summary, ensure_ascii=False, indent=2))
    print(f"results: {result_path}")
    print(f"summary: {summary_path}")
    print(f"excel: {excel_path}")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Evaluate current RAG QA accuracy for company demos."
    )
    subparsers = parser.add_subparsers(dest="command", required=True)

    inventory = subparsers.add_parser("inventory", help="Show upload and database readiness.")
    inventory.add_argument("--upload-root", default=str(DEFAULT_UPLOAD_ROOT), help="Upload directory to inspect.")
    inventory.add_argument("--out", default="", help="Optional JSON output path.")

    make_dataset = subparsers.add_parser("make-dataset", help="Generate QA test set from upload files or DB docs.")
    make_dataset.add_argument("--source", choices=["upload", "db"], default="upload")
    make_dataset.add_argument("--upload-root", default=str(DEFAULT_UPLOAD_ROOT))
    make_dataset.add_argument("--out", default=str(PROJECT_ROOT / "evaluate" / "company_demo_questions.json"))
    make_dataset.add_argument("--questions-per-doc", type=int, default=2)
    make_dataset.add_argument("--limit-docs", type=int, default=0, help="0 means no limit.")
    make_dataset.add_argument("--max-chars", type=int, default=12000)
    make_dataset.add_argument("--extensions", default=",".join(sorted(SUPPORTED_TEXT_EXTENSIONS)))
    make_dataset.add_argument("--include-unvectorized", action="store_true", help="Only applies to --source db.")
    make_dataset.add_argument("--base-url", default="")
    make_dataset.add_argument("--api-key", default="")
    make_dataset.add_argument("--model", default="")
    make_dataset.add_argument("--max-tokens", type=int, default=1600)
    make_dataset.add_argument("--temperature", type=float, default=0.2)

    make_image_dataset = subparsers.add_parser("make-image-dataset", help="Generate image-query QA test set from DB docs with images.")
    make_image_dataset.add_argument("--out", default=str(PROJECT_ROOT / "evaluate" / "image_retrieval_questions.json"))
    make_image_dataset.add_argument("--image-dir", default=str(PROJECT_ROOT / "evaluate" / "image_retrieval_images"))
    make_image_dataset.add_argument("--questions-per-doc", type=int, default=1)
    make_image_dataset.add_argument("--limit-docs", type=int, default=0, help="0 means no limit.")
    make_image_dataset.add_argument("--include-unvectorized", action="store_true")
    make_image_dataset.add_argument("--max-images-per-case", type=int, default=1)
    make_image_dataset.add_argument("--base-url", default="")
    make_image_dataset.add_argument("--api-key", default="")
    make_image_dataset.add_argument("--model", default="")
    make_image_dataset.add_argument("--max-tokens", type=int, default=1600)
    make_image_dataset.add_argument("--temperature", type=float, default=0.2)

    retrieve = subparsers.add_parser("retrieve", help="Evaluate retrieval accuracy only.")
    retrieve.add_argument("--dataset", required=True)
    retrieve.add_argument("--out-dir", default="")
    retrieve.add_argument("--limit-cases", type=int, default=0, help="0 means no limit.")
    retrieve.add_argument("--show-rag-debug", action="store_true")
    retrieve.add_argument("--use-source-images", action="store_true", help="Use source document images as query images when dataset has no query_images.")
    retrieve.add_argument("--max-query-images", type=int, default=1)

    run = subparsers.add_parser("run", help="Ask current RAG and judge accuracy.")
    run.add_argument("--dataset", required=True)
    run.add_argument("--out-dir", default="")
    run.add_argument("--limit-cases", type=int, default=0, help="0 means no limit.")
    run.add_argument("--api-key", default="")
    run.add_argument("--answer-base-url", default="")
    run.add_argument("--answer-model", default="")
    run.add_argument("--max-answer-tokens", type=int, default=2000)
    run.add_argument("--judge-base-url", default="")
    run.add_argument("--judge-model", default="")
    run.add_argument("--max-judge-tokens", type=int, default=600)
    run.add_argument("--judge-threshold", type=float, default=0.7)
    run.add_argument("--no-judge", action="store_true")
    run.add_argument("--show-rag-debug", action="store_true")
    run.add_argument("--use-source-images", action="store_true", help="Use source document images as query images when dataset has no query_images.")
    run.add_argument("--max-query-images", type=int, default=1)

    retrieve_api = subparsers.add_parser("retrieve-api", help="Evaluate retrieval through existing /message/ask API.")
    retrieve_api.add_argument("--dataset", required=True)
    retrieve_api.add_argument("--backend-url", default="http://127.0.0.1:8000")
    retrieve_api.add_argument("--out-dir", default="")
    retrieve_api.add_argument("--limit-cases", type=int, default=0, help="0 means no limit.")
    retrieve_api.add_argument("--token", default="")
    retrieve_api.add_argument("--username", default="")
    retrieve_api.add_argument("--password", default="")
    retrieve_api.add_argument("--role", default="")
    retrieve_api.add_argument("--use-source-images", action="store_true", help="Use source document images as query images when dataset has no query_images.")
    retrieve_api.add_argument("--max-query-images", type=int, default=1)

    run_api = subparsers.add_parser("run-api", help="Evaluate through existing /message/ask backend API.")
    run_api.add_argument("--dataset", required=True)
    run_api.add_argument("--backend-url", default="http://127.0.0.1:8000")
    run_api.add_argument("--out-dir", default="")
    run_api.add_argument("--limit-cases", type=int, default=0, help="0 means no limit.")
    run_api.add_argument("--token", default="")
    run_api.add_argument("--username", default="")
    run_api.add_argument("--password", default="")
    run_api.add_argument("--role", default="")
    run_api.add_argument("--api-key", default="")
    run_api.add_argument("--judge-base-url", default="")
    run_api.add_argument("--judge-model", default="")
    run_api.add_argument("--max-judge-tokens", type=int, default=600)
    run_api.add_argument("--judge-threshold", type=float, default=0.7)
    run_api.add_argument("--no-judge", action="store_true")
    run_api.add_argument("--use-source-images", action="store_true", help="Use source document images as query images when dataset has no query_images.")
    run_api.add_argument("--max-query-images", type=int, default=1)

    return parser


async def async_main() -> None:
    load_dotenv(PROJECT_ROOT / ".env")
    parser = build_parser()
    args = parser.parse_args()
    if hasattr(args, "limit_docs") and args.limit_docs == 0:
        args.limit_docs = None
    if hasattr(args, "limit_cases") and args.limit_cases == 0:
        args.limit_cases = None

    if args.command == "inventory":
        await command_inventory(args)
    elif args.command == "make-dataset":
        await command_make_dataset(args)
    elif args.command == "make-image-dataset":
        await command_make_image_dataset(args)
    elif args.command == "retrieve":
        await command_retrieve(args)
    elif args.command == "run":
        await command_run(args)
    elif args.command == "retrieve-api":
        await command_retrieve_api(args)
    elif args.command == "run-api":
        await command_run_api(args)
    else:
        raise ValueError(f"Unknown command: {args.command}")


if __name__ == "__main__":
    asyncio.run(async_main())

